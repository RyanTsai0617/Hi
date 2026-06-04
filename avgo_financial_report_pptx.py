"""AVGO（Broadcom）財報分析 — PPTX 產生器。

沿用 repo 既有視覺規範（etf_weekly_friday_preview_pptx.py）：
16:9、白底、teal #1B4F5A、微軟正黑體、yellow accent、footer 法人業務二處。

涵蓋：
- FY2025 全年 + Q1 FY2026 最新季報（資料來源見最後一頁，皆取自 SEC 8-K / 官方新聞稿）
- 營收 / AI 營收成長、segment 拆解、獲利與現金流、guidance + 投資觀點

本檔重用 etf 模組的 helper 與色彩常數，僅新增 AVGO 專用的 matplotlib 圖表與 slide builder。
matplotlib 中文以 WenQuanYi Zen Hei（Linux）/ Microsoft JhengHei（Windows）渲染；
PPTX 內文字以微軟正黑體 metadata 標記，PowerPoint 開檔時套用。

跑法： python3 avgo_financial_report_pptx.py  →  輸出 avgo_financial_report.pptx
"""
from __future__ import annotations
from pathlib import Path
from io import BytesIO

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 重用既有視覺 helper 與色彩常數（DRY，與週報同一套規範）
from etf_weekly_friday_preview_pptx import (
    apply_font, add_text, add_rect, add_rounded_rect, add_image,
    add_footer, add_title_bar, build_slide_section_divider,
    ACCENT, ACCENT_LIGHT, BODY, SUB, META, GAIN, LOSS, WHITE,
    LIGHT_BG, ALT_ROW, DIVIDER_YELLOW,
    CARD_BG_1, CARD_BG_2, CARD_BG_3,
    ACCENT_HEX, ACCENT_LIGHT_HEX, GAIN_HEX, LOSS_HEX, LIGHT_BG_HEX, META_HEX,
)

YELLOW_HEX = "#E8C547"

# matplotlib 中文字型：Linux 用 WenQuanYi Zen Hei，沒有就 fallback
for _fp in ["/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
            "C:/Windows/Fonts/msjh.ttc", "C:/Windows/Fonts/msjh.ttf"]:
    if Path(_fp).exists():
        font_manager.fontManager.addfont(_fp)
        break
plt.rcParams["font.family"] = ["WenQuanYi Zen Hei", "Microsoft JhengHei", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False


# ============================================================
#  AVGO 資料（全部取自官方新聞稿 / SEC 8-K，單位：US$ 百萬，除非另註）
#  FY 結束日：FY2025 = 2025-11-02、Q1 FY2026 = 約 2026-02-01
# ============================================================

# 近 5 季合併營收（US$ M）與 YoY
QUARTERS = [
    # (label, revenue, ai_revenue, yoy_pct)
    ("Q1'25", 14916, 4100, 25),
    ("Q2'25", 15004, 4400, 20),
    ("Q3'25", 15952, 5200, 22),
    ("Q4'25", 18015, 6300, 28),
    ("Q1'26", 19311, 8400, 29),
]

# 全年（FY2024 vs FY2025）
FY_REVENUE = [("FY2023", 35819), ("FY2024", 51574), ("FY2025", 62887)]
FY_NET_INCOME_GAAP = {"FY2024": 5895, "FY2025": 23126}
FY_FCF = {"FY2024": 19405, "FY2025": 26900}

# FY2025 segment（約略，US$ B）
SEGMENTS_FY25 = [("半導體解決方案\nSemiconductor", 37.0), ("基礎架構軟體\nInfrastructure SW", 26.0)]


# ============================================================
#  Matplotlib 圖表
# ============================================================

def _style_ax(ax):
    ax.set_facecolor(LIGHT_BG_HEX)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")
    ax.tick_params(colors=META_HEX)
    ax.grid(axis="y", linestyle=":", alpha=0.3)


def make_revenue_chart(*, width_in=11.5, height_in=4.4):
    """季營收長條（teal）＋ YoY 折線（yellow，右軸）。"""
    labels = [q[0] for q in QUARTERS]
    rev = [q[1] / 1000.0 for q in QUARTERS]  # → US$ B
    yoy = [q[3] for q in QUARTERS]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=150)
    fig.patch.set_facecolor("#FFFFFF")
    bars = ax.bar(labels, rev, color=ACCENT_HEX, width=0.6, edgecolor="none", zorder=3)
    for b, v in zip(bars, rev):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.3, f"${v:.1f}B",
                ha="center", va="bottom", fontsize=10, color=ACCENT_HEX, weight="bold")
    ax.set_ylabel("合併營收 (US$ B)", fontsize=10, color=META_HEX)
    ax.set_ylim(0, max(rev) * 1.18)
    _style_ax(ax)

    ax2 = ax.twinx()
    ax2.plot(labels, yoy, color=YELLOW_HEX, marker="o", markersize=7,
             linewidth=2.5, zorder=4, markeredgecolor="#B8860B")
    for x, y in zip(labels, yoy):
        ax2.text(x, y + 1.2, f"+{y}%", ha="center", fontsize=9, color="#8a6d00", weight="bold")
    ax2.set_ylabel("YoY 成長 (%)", fontsize=10, color="#8a6d00")
    ax2.set_ylim(0, max(yoy) * 1.6)
    ax2.spines["top"].set_visible(False)
    ax2.tick_params(colors="#8a6d00")

    ax.set_title("合併營收 — 連續創高、YoY 加速至 +29%", fontsize=13,
                 color=ACCENT_HEX, pad=10, weight="bold")
    plt.tight_layout()
    return _save(fig)


def make_ai_revenue_chart(*, width_in=11.5, height_in=4.4):
    """AI 半導體營收成長（季）＋ 佔比文字。"""
    labels = [q[0] for q in QUARTERS]
    ai = [q[2] / 1000.0 for q in QUARTERS]  # → US$ B
    share = [q[2] / q[1] * 100 for q in QUARTERS]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=150)
    fig.patch.set_facecolor("#FFFFFF")
    bars = ax.bar(labels, ai, color=ACCENT_LIGHT_HEX, width=0.6, edgecolor="none", zorder=3)
    for b, v, s in zip(bars, ai, share):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.15, f"${v:.1f}B",
                ha="center", va="bottom", fontsize=10, color=ACCENT_LIGHT_HEX, weight="bold")
        ax.text(b.get_x() + b.get_width() / 2, v / 2, f"佔總營收\n{s:.0f}%",
                ha="center", va="center", fontsize=8.5, color="#FFFFFF")
    ax.set_ylabel("AI 半導體營收 (US$ B)", fontsize=10, color=META_HEX)
    ax.set_ylim(0, max(ai) * 1.2)
    _style_ax(ax)
    ax.set_title("AI 半導體營收 — Q1'26 達 $8.4B、YoY +106%", fontsize=13,
                 color=ACCENT_HEX, pad=10, weight="bold")
    plt.tight_layout()
    return _save(fig)


def make_segment_donut(*, width_in=5.4, height_in=4.4):
    """FY2025 segment donut（圖例置於下方，避免標籤被裁切）。"""
    labels = [s[0].replace("\n", " ") for s in SEGMENTS_FY25]
    vals = [s[1] for s in SEGMENTS_FY25]
    colors = [ACCENT_HEX, YELLOW_HEX]
    share = [v / sum(vals) * 100 for v in vals]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=150)
    fig.patch.set_facecolor("#FFFFFF")
    wedges, _ = ax.pie(vals, colors=colors, startangle=90, counterclock=False,
                       wedgeprops=dict(width=0.42, edgecolor="#FFFFFF", linewidth=2),
                       autopct=None)
    ax.text(0, 0.12, "$62.9B", ha="center", va="center",
            fontsize=22, color=ACCENT_HEX, weight="bold")
    ax.text(0, -0.24, "FY2025 總營收", ha="center", va="center",
            fontsize=10, color=META_HEX)
    ax.set(aspect="equal")
    legend_labels = [f"{lab}  ${v:.0f}B · {s:.0f}%" for lab, v, s in zip(labels, vals, share)]
    ax.legend(wedges, legend_labels, loc="upper center", bbox_to_anchor=(0.5, 0.04),
              ncol=1, frameon=False, fontsize=9, handlelength=1.1, labelspacing=0.6)
    plt.subplots_adjust(top=0.98, bottom=0.16)
    return _save(fig)


def make_fcf_chart(*, width_in=6.2, height_in=4.4):
    """FY 自由現金流 vs GAAP 淨利。"""
    years = ["FY2024", "FY2025"]
    fcf = [FY_FCF[y] / 1000.0 for y in years]
    ni = [FY_NET_INCOME_GAAP[y] / 1000.0 for y in years]
    x = range(len(years))
    w = 0.36

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=150)
    fig.patch.set_facecolor("#FFFFFF")
    b1 = ax.bar([i - w / 2 for i in x], fcf, w, color=ACCENT_HEX, label="自由現金流 FCF", zorder=3)
    b2 = ax.bar([i + w / 2 for i in x], ni, w, color=CARD_BG_1_hex(), label="GAAP 淨利", zorder=3)
    for bars in (b1, b2):
        for b in bars:
            ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.4,
                    f"${b.get_height():.1f}B", ha="center", va="bottom",
                    fontsize=9, color=ACCENT_HEX, weight="bold")
    ax.set_xticks(list(x))
    ax.set_xticklabels(years)
    ax.set_ylabel("US$ B", fontsize=10, color=META_HEX)
    ax.set_ylim(0, max(fcf + ni) * 1.2)
    _style_ax(ax)
    ax.legend(fontsize=9, loc="upper left", frameon=False)
    ax.set_title("現金流與獲利躍升", fontsize=13, color=ACCENT_HEX, pad=10, weight="bold")
    plt.tight_layout()
    return _save(fig)


def CARD_BG_1_hex():
    return YELLOW_HEX


def _save(fig):
    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="#FFFFFF", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


# ============================================================
#  Slide builders
# ============================================================

def build_cover(slide, footer):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), ACCENT)
    add_rect(slide, Inches(0), Inches(0.4), Inches(13.333), Inches(0.05), DIVIDER_YELLOW)

    add_text(slide, Inches(0.6), Inches(1.15), Inches(12.1), Inches(1.0),
             "Broadcom（AVGO）財報分析",
             size=36, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.6), Inches(2.25), Inches(12.1), Inches(0.6),
             "FY2025 全年 · Q1 FY2026 最新季報（截至 2026/02/01）",
             size=17, color=BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.6), Inches(2.85), Inches(12.1), Inches(0.5),
             "NASDAQ: AVGO · 半導體 + 基礎架構軟體 · AI 客製化加速器（XPU）與 AI 網通領導者",
             size=12, color=SUB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cards = [
        ("FY2025 營收", "$62.9B", "+22% YoY", CARD_BG_3),
        ("FY2025 AI 營收", "$20B", "Q1'26 季增至 $8.4B", CARD_BG_2),
        ("FY2025 自由現金流", "$26.9B", "FCF margin 約 43%", CARD_BG_1),
    ]
    card_w, card_h, gap = Inches(3.5), Inches(2.25), Inches(0.3)
    total_w = card_w * 3 + gap * 2
    start_x = (Inches(13.333) - total_w) / 2
    cards_y = Inches(3.85)
    for i, (label, value, sub, bg) in enumerate(cards):
        cx = start_x + i * (card_w + gap)
        add_rounded_rect(slide, cx, cards_y, card_w, card_h, bg)
        add_text(slide, cx, cards_y + Inches(0.22), card_w, Inches(0.45), label,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx, cards_y + Inches(0.72), card_w, Inches(0.9), value,
                 size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx, cards_y + Inches(1.68), card_w, Inches(0.45), sub,
                 size=11, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.4),
             "資料來源：Broadcom 官方新聞稿 / SEC 8-K（FY2025 Q4 & Q1 FY2026）",
             size=10, color=META, align=PP_ALIGN.CENTER)
    _dark_footer(slide, footer)


def build_q1_highlights(slide, footer):
    add_title_bar(slide, "Q1 FY2026 重點", "截至 2026/02/01、2026/03/04 公布 · 全面優於財測")

    kpis = [
        ("合併營收", "$19.31B", "+29% YoY", GAIN),
        ("AI 半導體營收", "$8.4B", "+106% YoY", GAIN),
        ("半導體營收（紀錄）", "$12.5B", "+52% YoY", GAIN),
        ("Non-GAAP 淨利", "$10.19B", "EPS $2.05", GAIN),
        ("營業利益率", "66.4%", "+50 bps YoY", GAIN),
        ("Adj. EBITDA", "$13.1B", "佔營收 68%", GAIN),
        ("自由現金流", "$8.01B", "佔營收 41%", GAIN),
        ("季配息", "$0.65", "每股 / 季", ACCENT),
    ]
    cols, rows = 4, 2
    gx, gy = Inches(0.6), Inches(1.75)
    cw, ch, gap = Inches(2.92), Inches(2.15), Inches(0.18)
    for idx, (label, value, sub, vcolor) in enumerate(kpis):
        r, c = divmod(idx, cols)
        x = gx + c * (cw + gap)
        y = gy + r * (ch + Inches(0.3))
        add_rounded_rect(slide, x, y, cw, ch, LIGHT_BG)
        add_rect(slide, x, y, cw, Inches(0.12), ACCENT)
        add_text(slide, x + Inches(0.15), y + Inches(0.3), cw - Inches(0.3), Inches(0.45),
                 label, size=12, bold=True, color=SUB, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.15), y + Inches(0.8), cw - Inches(0.3), Inches(0.75),
                 value, size=30, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.15), y + Inches(1.55), cw - Inches(0.3), Inches(0.45),
                 sub, size=12, bold=True, color=vcolor, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.45),
             "重點：AI 加速器（XPU）＋ AI 網通雙引擎，帶動半導體單季創高；軟體（VMware）穩定貢獻現金流。",
             size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 3, TOTAL, footer)


def build_revenue_slide(slide, footer):
    add_title_bar(slide, "營收趨勢", "近 5 季合併營收 + YoY 成長率")
    add_image(slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.4), make_revenue_chart())
    add_text(slide, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.6),
             "FY2025 全年營收 $62.9B（+22% YoY，FY2024 $51.6B），Q1'26 單季 $19.3B 再創新高、YoY 連續加速。",
             size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 4, TOTAL, footer)


def build_ai_slide(slide, footer):
    add_title_bar(slide, "AI 營收引擎", "客製化加速器（XPU）+ AI 網通（Ethernet switch）")
    add_image(slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.4), make_ai_revenue_chart())
    add_text(slide, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.6),
             "AI 營收佔比由 Q1'25 約 27% 拉升至 Q1'26 約 44%；管理層看 2027 年 AI 晶片營收上看 $100B+、AI backlog 約 $73B（18 個月交期）。",
             size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 5, TOTAL, footer)


def build_segment_cashflow_slide(slide, footer):
    add_title_bar(slide, "業務結構與現金流", "FY2025 雙引擎：半導體 + 基礎架構軟體")
    add_image(slide, Inches(0.5), Inches(1.75), Inches(5.4), Inches(4.4), make_segment_donut())
    add_image(slide, Inches(6.3), Inches(1.75), Inches(6.2), Inches(4.4), make_fcf_chart())
    add_text(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5),
             "軟體（VMware 整併）提供高毛利、可預期的訂閱現金流；半導體（AI）提供高成長。GAAP 淨利由 FY2024 $5.9B 躍升至 FY2025 $23.1B。",
             size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 6, TOTAL, footer)


def build_outlook_slide(slide, footer):
    add_title_bar(slide, "展望與投資觀點", "Q2 FY2026 財測 + 中長期 AI 敘事 + 對台灣機構意義")

    # 左欄：財測與里程碑
    add_text(slide, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.4),
             "📈 財測與里程碑", size=15, bold=True, color=ACCENT_LIGHT)
    left_items = [
        "Q2 FY2026 營收財測約 $22.0B（+47% YoY）",
        "Q2 AI 半導體營收預估約 $10.7B",
        "2027 年 AI 晶片營收 line of sight 上看 $100B+",
        "AI backlog 約 $73B、18 個月交期能見度高",
        "毛利／營益率維持高檔（營益率 66%+）",
        "強勁 FCF 支撐配息 + 去槓桿（VMware 併購債）",
    ]
    add_text(slide, Inches(0.6), Inches(2.2), Inches(6.0), Inches(4.0),
             "\n".join(f"•  {t}" for t in left_items), size=12, color=BODY)

    # 右欄：對台灣機構意義
    add_text(slide, Inches(7.0), Inches(1.7), Inches(5.7), Inches(0.4),
             "🇹🇼 對台灣機構意義", size=15, bold=True, color=ACCENT_LIGHT)
    right_items = [
        "AVGO 為 TSMC 先進製程（CoWoS / N3 / N2）核心大客戶之一",
        "XPU 放量 → 台系 ABF 載板、CCL、測試、散熱供應鏈受惠",
        "AI 網通（Tomahawk / Jericho）利多交換器、光通訊族群",
        "可作為 SMH / SOXX / AIQ 等 AI 半導體 ETF 的權重觀察錨",
        "風險：客製化 ASIC 客戶集中度、軟體成長趨緩、評價偏高",
    ]
    add_text(slide, Inches(7.0), Inches(2.2), Inches(5.7), Inches(4.0),
             "\n".join(f"•  {t}" for t in right_items), size=12, color=BODY)

    add_footer(slide, 8, TOTAL, footer)


def build_disclaimer_slide(slide, footer):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), ACCENT)
    add_text(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.7),
             "資料來源與聲明", size=26, bold=True, color=WHITE)
    add_rect(slide, Inches(0.8), Inches(1.55), Inches(1.2), Inches(0.08), DIVIDER_YELLOW)

    src = [
        "Broadcom Inc.「Q1 FY2026 Financial Results」官方新聞稿（2026/03/04）",
        "Broadcom Inc.「Q4 & FY2025 Financial Results」官方新聞稿（2025/12/11）",
        "SEC Form 8-K 財報附件（CIK 0001730168，FY2025 / FY2026 各季）",
        "全年數字部分由各季 SEC 8-K 加總推算；segment 為約略值，以官方 10-K 為準。",
    ]
    add_text(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.4),
             "\n".join(f"•  {t}" for t in src), size=13, color=WHITE)

    add_text(slide, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.8),
             "本簡報僅為財報數據彙整與分析，供內部研究與法人業務參考，"
             "不構成任何投資建議或要約。投資人應自行評估風險。\n"
             "金額單位為美元（US$），會計期間以 Broadcom 財報年度（FY 結束於 11 月初）為準。",
             size=12, color=RGBColor(0xD5, 0xE5, 0xE8))

    add_text(slide, Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3),
             footer, size=9, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _dark_footer(slide, footer):
    add_text(slide, Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3),
             footer, size=9, color=META, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
#  Main
# ============================================================

TOTAL = 9


def build(out_path="avgo_financial_report.pptx"):
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    blank = pres.slide_layouts[6]
    footer = "國泰證券法人債券業務部 · 法人業務二處 · AVGO 財報分析"

    build_cover(pres.slides.add_slide(blank), footer)
    build_slide_section_divider(pres.slides.add_slide(blank), "01", "營運總覽 · Financial Overview", footer)
    build_q1_highlights(pres.slides.add_slide(blank), footer)
    build_revenue_slide(pres.slides.add_slide(blank), footer)
    build_ai_slide(pres.slides.add_slide(blank), footer)
    build_segment_cashflow_slide(pres.slides.add_slide(blank), footer)
    build_slide_section_divider(pres.slides.add_slide(blank), "02", "展望 · Outlook & Thesis", footer)
    build_outlook_slide(pres.slides.add_slide(blank), footer)
    build_disclaimer_slide(pres.slides.add_slide(blank), footer)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    pres.save(out_path)
    print(f"[OK] AVGO financial report saved: {out_path}  ({len(pres.slides._sldIdLst)} slides)")
    return out_path


if __name__ == "__main__":
    build()
