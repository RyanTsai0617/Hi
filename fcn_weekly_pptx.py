"""FCN Weekly — KI/KO 期望值週報 PPTX builder (GS style).

沿用 etf_weekly_friday_preview_pptx_v3_gs.py 的 Goldman 視覺（navy + gold、
直角、高密度表格、BOTTOM LINE callout），8 slides 16:9。
純 python-pptx shapes/tables，不用 matplotlib（遠端 Linux 無中文字型，
圖片會豆腐字；pptx 文字字型由開啟端渲染，不受影響）。

數據對應 _draft_fcn_weekly_2026-07-03.md（20260706 期，est. 值 Monday 校正）。
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============== Goldman Sachs visual palette ==============
GS_BLUE = RGBColor(0x00, 0x29, 0x5C)
GS_BLUE_LIGHT = RGBColor(0x1F, 0x4E, 0x82)
GS_GOLD = RGBColor(0xFF, 0xB8, 0x1C)
GS_GOLD_DARK = RGBColor(0xC8, 0x8E, 0x10)
BLACK = RGBColor(0x10, 0x10, 0x10)
DARK_GRAY = RGBColor(0x2A, 0x2A, 0x2A)
MED_GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
ALT_ROW = RGBColor(0xF5, 0xF5, 0xF5)
BORDER_GRAY = RGBColor(0xC8, 0xC8, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GAIN = RGBColor(0x1A, 0x6E, 0x1A)
LOSS = RGBColor(0xA5, 0x0B, 0x0B)

ZH_FONT = "微軟正黑體"
EN_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ISSUE = "20260706"
GEN_DATE = "2026-07-03"
WEEK_RANGE = "2026-06-29 ~ 2026-07-02"


# ============== primitives ==============

def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, wrap=True, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing != 1.0:
            p.line_spacing = line_spacing
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = ZH_FONT
    return tb


def add_banner(slide, page_no, section="FCN WEEKLY | KI/KO EXPECTED VALUE"):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.35), GS_BLUE)
    add_text(slide, Inches(0.35), Inches(0.04), Inches(9), Inches(0.28),
             [[(f"STRUCTURED PRODUCTS  |  {section}", 11, True, WHITE)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(12.3), Inches(0.04), Inches(0.7), Inches(0.28),
             [[(str(page_no), 11, True, GS_GOLD)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_table(slide, x, y, w, h, headers, rows, col_widths=None,
              font_size=10.5, header_size=10.5, align_map=None,
              color_cols=None):
    """color_cols: set of column indexes whose leading +/- decides GAIN/LOSS."""
    shp = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    tbl = shp.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for j, htext in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GS_BLUE
        cell.margin_left = cell.margin_right = Inches(0.05)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        p.alignment = (align_map or {}).get(j, PP_ALIGN.LEFT)
        r = p.add_run()
        r.text = htext
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = ZH_FONT
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ALT_ROW if i % 2 else WHITE
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            p.alignment = (align_map or {}).get(j, PP_ALIGN.LEFT)
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            r.font.name = ZH_FONT
            color = DARK_GRAY
            if color_cols and j in color_cols:
                s = str(val)
                if s.startswith("+"):
                    color = GAIN
                elif s.startswith("-"):
                    color = LOSS
            r.font.color.rgb = color
            r.font.bold = (j == 0)
    return tbl


def add_bottom_line(slide, x, y, w, text):
    add_rect(slide, x, y, w, Pt(2.5), GS_GOLD)
    add_text(slide, x, y + Pt(5), w, Inches(0.55),
             [[("BOTTOM LINE  ", 10.5, True, GS_GOLD_DARK),
               (text, 10.5, True, DARK_GRAY)]], line_spacing=1.05)


# ============== slides ==============

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, 1)
    add_text(slide, Inches(0.7), Inches(0.9), Inches(12), Inches(1.5),
             [[("FCN Weekly — KI/KO 期望值週報", 40, True, BLACK)],
              [(f"{ISSUE} 期  |  週次 {WEEK_RANGE}  |  生成 {GEN_DATE}（Fri）", 15, False, MED_GRAY)]],
             line_spacing=1.1)
    add_text(slide, Inches(0.7), Inches(2.35), Inches(12), Inches(0.9),
             [[("E[FCN] = p(KO)×小獲利 + p(平安)×票息 + p(KI)×大損失", 18, True, GS_BLUE)],
              [("發生機率與發生損失的差異，造成期望值的不同 — 存量健檢 × 發行時點雙主軸", 13, False, DARK_GRAY)]],
             line_spacing=1.25)

    cards = [
        (GS_BLUE, WHITE, "VIX 16.59", "殺盤週衝高後回落", "新發行票息原料流失中"),
        (GS_GOLD, BLACK, "非農 +57K", "vs 預期 ~110K", "9月升息機率 67% → <50%"),
        (DARK_GRAY, WHITE, "MU 高點 -13%", "SOX 高點回落 -6.7%", "存續 FCN 之 KI 緩衝被吃掉"),
    ]
    cx = Inches(0.7)
    for bg, fg, big, mid, small in cards:
        add_rect(slide, cx, Inches(3.6), Inches(3.9), Inches(2.0), bg)
        add_text(slide, cx + Inches(0.25), Inches(3.8), Inches(3.4), Inches(1.7),
                 [[(big, 24, True, fg)], [(mid, 12, False, fg)],
                  [(small, 11, False, fg)]], line_spacing=1.3)
        cx += Inches(4.15)

    add_text(slide, Inches(0.7), Inches(6.15), Inches(12), Inches(0.8),
             [[("⚠ v0 partial scan：yfinance 遭 proxy 封鎖，est. 數字以 WebSearch 新聞推算，Monday 校正。", 11, False, MED_GRAY)],
              [("本週僅 4 個交易日（7/3 休市、7/4 週六）→ Mon-Thu close 即完整週。", 11, False, MED_GRAY)]],
             line_spacing=1.3)
    add_text(slide, Inches(0.7), Inches(7.0), Inches(12), Inches(0.35),
             [[("國泰證券法人債券業務部 · 法人業務二處", 11, True, GS_BLUE)]])


def slide_framework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, 2)
    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.5),
             [[("〇  本期框架：KI/KO 期望值分解", 24, True, BLACK)]])
    add_text(slide, Inches(0.7), Inches(1.25), Inches(12), Inches(0.45),
             [[("E[FCN] = p(KO) × (已收票息+提前出場) + p(平安到期) × (全期票息) + p(KI) × (已收票息 − 承接跌幅)",
                14, True, GS_BLUE)]])

    headers = ["維度", "KO（敲出）", "KI（敲入）"]
    rows = [
        ["發生機率", "高：障礙設在 spot 附近，正常波動即觸發", "低：障礙設在深度價外（常見 60~70%）"],
        ["發生損益", "小額封頂（票息出場）", "大額開放（跟跌現貨，理論可至 -100%）"],
        ["期望值貢獻", "高機率 × 小獲利 → 穩定正貢獻", "低機率 × 大損失 → 不穩定大額負貢獻"],
        ["參數敏感度", "低（歷史波動即可估好）", "高：尾部機率最易被低估，期望值的主戰場"],
    ]
    add_table(slide, Inches(0.7), Inches(1.85), Inches(11.9), Inches(2.6),
              headers, rows,
              col_widths=[Inches(1.7), Inches(5.1), Inches(5.1)], font_size=11.5,
              header_size=12)

    add_text(slide, Inches(0.7), Inches(4.75), Inches(11.9), Inches(1.5),
             [[("為什麼是本週：", 12.5, True, BLACK),
               ("6/22-26 半導體殺盤把存續 FCN 的 KI 緩衝一次吃掉一大截 — 期望值惡化不是票息變了，是 p(KI) 變了；",
                12.5, False, DARK_GRAY)],
              [("同時 VIX 回落至 16.59，新發行的票息原料（隱含波動）正在流失。存量與增量兩端都踩在框架要害上。",
                12.5, False, DARK_GRAY)]], line_spacing=1.35)
    add_bottom_line(slide, Inches(0.7), Inches(6.35), Inches(11.9),
                    "機率估錯 2%（3% 估成 1%）× 損失 -40% = 期望值差 0.8pp — 對票息 6% 商品是重大差異。")


def slide_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, 3)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.5),
             [[("一  本週市場總覽（FCN 視角）", 24, True, BLACK)]])

    add_text(slide, Inches(0.7), Inches(1.2), Inches(5.6), Inches(0.35),
             [[("波動率 — 票息的原料", 13, True, GS_BLUE)]])
    vol_rows = [
        ["VIX", "16.59", "殺盤週衝高後回落，票息壓縮中"],
        ["9/18 升息機率", "<50%", "非農後自 67% 退潮"],
        ["美 10Y", "~4.46%", "FCN vs 直債比較基準"],
    ]
    add_table(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(1.5),
              ["指標", "本週", "FCN 意涵"], vol_rows,
              col_widths=[Inches(1.55), Inches(1.05), Inches(3.0)], font_size=10)

    add_text(slide, Inches(0.7), Inches(3.4), Inches(5.6), Inches(3.5),
             [[("Macro 事件", 13, True, GS_BLUE)],
              [("1. 6月非農 +57K vs ~110K，四個月最低；失業率 4.2%", 10.5, False, DARK_GRAY)],
              [("2. 週一 SMH +3%、NDX +2.3% 技術性修復；Broadcom 指引下修餘波未平", 10.5, False, DARK_GRAY)],
              [("3. Dow 首破 52,000；防禦板塊週漲 2%+；GOOGL 納入道瓊 +4.6%", 10.5, False, DARK_GRAY)],
              [("4. 金價自 8 個月低點反彈站回 $4,100（升息預期退潮驅動）", 10.5, False, DARK_GRAY)],
              [("5. 7/3 休市（7/4 週六）→ 本週 4 交易日即完整週", 10.5, False, DARK_GRAY)]],
             line_spacing=1.45)

    add_text(slide, Inches(6.7), Inches(1.2), Inches(5.9), Inches(0.35),
             [[("標的資產週表現（est. — Monday 校正）", 13, True, GS_BLUE)]])
    px_rows = [
        ["SMH", "+1.8", "-6.50", "週一 +3% 反彈後動能衰減"],
        ["SOXX", "+1.5", "-6.20", "Q2 累漲 ~94% 後高位震盪"],
        ["QQQ", "+1.2", "-3.20", "週四非農後 -0.8%"],
        ["TSLA", "+6.0", "n/a", "週一單日 +8.06% 領漲 Mag7"],
        ["MU", "-1.0", "-8.50", "記憶體供給過剩疑慮未消"],
        ["NVDA", "+2.5", "-3.00", "反彈但量能不足"],
        ["GLD", "+2.5", "-4.80", "站回 $4,100"],
        ["TLT", "+1.5", "+2.10", "非農爆冷、殖利率回落"],
        ["XLU/XLP", "+2.0", "+1.4", "防禦輪動、Dow 52,000"],
    ]
    add_table(slide, Inches(6.7), Inches(1.6), Inches(5.9), Inches(4.6),
              ["標的", "本週%", "上週%", "敘事"], px_rows,
              col_widths=[Inches(1.05), Inches(0.85), Inches(0.85), Inches(3.15)],
              font_size=10,
              align_map={1: PP_ALIGN.RIGHT, 2: PP_ALIGN.RIGHT},
              color_cols={1, 2})

    add_bottom_line(slide, Inches(0.7), Inches(6.6), Inches(11.9),
                    "反彈週 + vol 回落：存量 KI 壓力暫緩，但新發行票息窗口也在關閉 — 兩者都指向本期雙主軸。")


def topic_card(slide, x, y, w, h, no, title, lines):
    add_rect(slide, x, y, w, h, WHITE, line=BORDER_GRAY, line_w=Pt(1))
    add_text(slide, x + Inches(0.18), y + Inches(0.12), Inches(0.75), Inches(0.5),
             [[(no, 24, True, GS_GOLD_DARK)]])
    add_text(slide, x + Inches(0.95), y + Inches(0.16), w - Inches(1.15), Inches(0.55),
             [[(title, 13.5, True, BLACK)]], line_spacing=1.0)
    add_text(slide, x + Inches(0.18), y + Inches(0.78), w - Inches(0.4), h - Inches(0.95),
             [[(t, 10.5, bold, c)] for t, bold, c in lines], line_spacing=1.25)


def slide_main_topics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, 4)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.5),
             [[("三  主議題候選（3 選 1）", 24, True, BLACK)]])
    cards = [
        ("01", "半導體 FCN 存量健檢 — 殺盤後 KI 緩衝重估",
         [("數據：MU 自高點 -13%、SOX -6.7% 後僅部分收復", False, DARK_GRAY),
          ("標的：NVDA / AMD / MU / TSM ADR / SMH", False, DARK_GRAY),
          ("殺盤前 spot 定價的存續 FCN，KI 距離被吃掉 1/3~1/2；p(KI) 上修直接拉低期望值", False, DARK_GRAY),
          ("意義：盤點 5-6 月半導體 vintage；距 KI <20% 列觀察名單，評估 restructure 或對沖", True, GS_BLUE)]),
        ("02", "VIX 16.59 的發行時點課 — 波動率與票息窗口",
         [("數據：VIX 由殺盤週高點回落至 16.59，indicative 票息同步壓縮", False, DARK_GRAY),
          ("標的：SMH / QQQ / NVDA 高 IV 族群", False, DARK_GRAY),
          ("FCN 票息本質是賣 put 權利金，IV 是原料；平靜期發行 = 低價賣保險", False, DARK_GRAY),
          ("意義：建立 vol 觸發式發行紀律（VIX>22 或 IV 一年 75 分位）", True, GS_BLUE)]),
        ("03", "非農爆冷 × 升息預期退潮 — 利率/股權商品蹺蹺板",
         [("數據：9月升息機率 67% → <50%；10Y ~4.46%；TLT +1.5% est.", False, DARK_GRAY),
          ("標的：TLT / LQD / AGG vs 股權掛鉤 FCN", False, DARK_GRAY),
          ("債券底層止穩、鎖息吸引力回升；equity vol 回落壓縮 FCN 票息", False, DARK_GRAY),
          ("意義：KI 厭惡客戶的「直債 vs FCN」推薦順位需要重排", True, GS_BLUE)]),
    ]
    cx = Inches(0.55)
    for no, title, lines in cards:
        topic_card(slide, cx, Inches(1.35), Inches(4.05), Inches(4.8), no, title, lines)
        cx += Inches(4.22)
    add_bottom_line(slide, Inches(0.7), Inches(6.5), Inches(11.9),
                    "建議主議題 01：存量健檢最貼近法人客戶的當下痛點，且直接展演 p(KI) 重估如何改寫期望值。")


def slide_supp_topics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, 5)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.5),
             [[("三  補充議題候選（3 選 2）", 24, True, BLACK)]])
    cards = [
        ("04", "黃金連結票據 — $4,100 反彈後 KI 緩衝重建",
         [("GLD +2.5% est.，金價自 8 個月低點反彈", False, DARK_GRAY),
          ("止穩使黃金 FCN 的 KI 緩衝開始重建；但高盛 $4,900 目標下修在前", False, DARK_GRAY),
          ("意義：已 KI 觀察名單的減壓評估；新發行等緩衝 >25% 再進", True, GS_BLUE)]),
        ("05", "防禦板塊 FCN — 低票息、低 p(KI) 的另一端",
         [("XLU/XLV/XLP 週漲 2%+，Dow 首破 52,000", False, DARK_GRAY),
          ("低票息 × 低 KI 機率的組合：severity 端優勢，KI 後條件損失小", False, DARK_GRAY),
          ("意義：KI 厭惡型客戶產品線；對照表說明「票息差距買的是什麼風險」", True, GS_BLUE)]),
        ("06", "TSLA 單週 +8% — 高 IV 標的教材案例",
         [("TSLA 週一 +8.06% 領漲 Mag7，雙向波動極大", False, DARK_GRAY),
          ("高 IV = 高票息 = 高 p(KI)：票息常為大盤標的 2-3 倍，但 KI 緩衝兩週可蒸發", False, DARK_GRAY),
          ("意義：客戶教育素材 — 票息不是報酬率，是尾部風險的對價", True, GS_BLUE)]),
    ]
    cx = Inches(0.55)
    for no, title, lines in cards:
        topic_card(slide, cx, Inches(1.35), Inches(4.05), Inches(4.3), no, title, lines)
        cx += Inches(4.22)
    add_bottom_line(slide, Inches(0.7), Inches(6.1), Inches(11.9),
                    "建議搭配 05 + 06：一端展示低 KI 的代價（薄票息），一端展示高票息的代價（尾部）— 框架首尾呼應。")


def slide_matrix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, 6)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.5),
             [[("四  FCN 條款試算矩陣", 24, True, BLACK),
               ("   illustrative · 非報價 · 6M / 月觀察 / KI 65% / KO 100%", 12, False, MED_GRAY)]])
    headers = ["標的", "現況定位", "IV", "示意年化票息", "KI 緩衝", "近月最大回撤", "期望值註記"]
    rows = [
        ["MU", "自高點 -13% 未收復", "極高", "~18-22%", "35%", "-13%", "一週吃掉 1/3 緩衝，p(KI) 全表最高"],
        ["NVDA", "反彈中，量能不足", "高", "~14-17%", "35%", "-8%", "票息肥但 AI 估值重設未完成"],
        ["TSLA", "週一 +8% 高波動雙向", "極高", "~18-24%", "35%", "-9%", "severity 與 frequency 雙高"],
        ["TSM ADR", "基本面最穩 Q1 +35% YoY", "中高", "~10-13%", "35%", "-6%", "半導體中 p(KI) 相對低的核心"],
        ["SMH", "ETF 分散單一股風險", "中高", "~9-12%", "35%", "-7%", "以 ETF 換單一股：severity 降檔"],
        ["QQQ", "指數級分散", "中", "~7-9%", "35%", "-5%", "票息薄但尾部最淺"],
        ["GLD", "8 個月低點反彈中", "中", "~6-8%", "35%", "-6%", "等緩衝重建 >25% 再進"],
        ["XLU", "防禦輪動受惠", "低", "~4-6%", "35%", "-3%", "低票息低 KI：KI 厭惡客戶適配"],
    ]
    add_table(slide, Inches(0.55), Inches(1.35), Inches(12.25), Inches(4.7),
              headers, rows,
              col_widths=[Inches(1.0), Inches(2.15), Inches(0.7), Inches(1.35),
                          Inches(0.9), Inches(1.25), Inches(4.9)],
              font_size=10.5, header_size=11,
              align_map={2: PP_ALIGN.CENTER, 3: PP_ALIGN.RIGHT,
                         4: PP_ALIGN.RIGHT, 5: PP_ALIGN.RIGHT},
              color_cols={5})
    add_bottom_line(slide, Inches(0.7), Inches(6.5), Inches(11.9),
                    "票息由高 IV 單一股向低 IV 指數/防禦遞減，p(KI) 與條件損失同向遞減 — 票息差距買的是尾部風險。")


def slide_eki_compare(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, 7)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.5),
             [[("四之二  同 K 對照 — 轉換機率 vs 預期損失", 24, True, BLACK),
               ("   Monte Carlo · σ=40% · 6M · EKI 65% 到期觀察 · 非年化", 12, False, MED_GRAY)]])
    headers = ["K", "轉換機率  無KI → EKI 65%", "降幅", "預期損失  無KI → EKI 65%", "降幅"]
    rows = [
        ["80%", "23.6% → 7.3%", "−69%", "2.69% → 1.60%", "−41%"],
        ["90%", "38.1% → 7.3%", "−81%", "5.70% → 2.32%", "−59%"],
        ["100%", "52.8% → 7.3%", "−86%", "10.16% → 3.04%", "−70%"],
    ]
    add_table(slide, Inches(0.55), Inches(1.35), Inches(7.3), Inches(1.9),
              headers, rows,
              col_widths=[Inches(0.75), Inches(2.5), Inches(0.9), Inches(2.5),
                          Inches(0.65)],
              font_size=12, header_size=10.5,
              align_map={0: PP_ALIGN.RIGHT, 1: PP_ALIGN.CENTER, 2: PP_ALIGN.RIGHT,
                         3: PP_ALIGN.CENTER, 4: PP_ALIGN.RIGHT},
              color_cols={2, 4})
    add_text(slide, Inches(0.55), Inches(3.55), Inches(7.3), Inches(2.8),
             [[("機率降幅永遠 > 預期損失降幅（差 16~28pp）", 13, True, GS_BLUE)],
              [("EKI 刪掉的是高頻低損的淺層轉換（65%<S<K）：「次數」掉得快、「錢」掉得慢；深層情境（S<65%）有無 EKI 損益完全相同 — 免小傷、不免大傷。",
                11, False, DARK_GRAY)],
              [("兩個降幅的差距 = severity 濃縮", 13, True, GS_BLUE)],
              [("每 1% 轉換機率攜帶的預期損失上升：K=100% 時 19.6% → 42.4%，最低損失自 ~0% 跳到 K−KI = 35%。",
                11, False, DARK_GRAY)]],
             line_spacing=1.3)
    add_text(slide, Inches(8.15), Inches(1.35), Inches(4.6), Inches(5.0),
             [[("轉換機率與 K 脫鉤", 12, True, GS_BLUE)],
              [("EKI 後三檔 K 都是 7.3% = P(S<KI)。K 只剩「決定轉換深度」— EKI 讓 K 可以積極的機制；但 K 80→100% 預期損失僅 +1.44pp、懸崖 +20pp，對客戶是純 severity", 10.5, False, DARK_GRAY)],
              [("溝通與比價各取一欄", 12, True, GS_BLUE)],
              [("「轉換機率 −86%」是賣點話術欄；「預期損失只 −70%、剩餘轉換深度翻倍」是完整敘述欄。預期損失欄 = 公平票息讓渡上限：票息差 < 預期損失差 → 客戶在補貼發行端", 10.5, False, DARK_GRAY)],
              [("觀察頻率是隱藏參數", 12, True, GS_BLUE)],
              [("KI 65% 改每日觀察觸發率 7.3%→~12.9%，且模型未含跳空（MU 單週 −13%）", 10.5, False, DARK_GRAY)]],
             line_spacing=1.25)
    add_bottom_line(slide, Inches(0.7), Inches(6.5), Inches(11.9),
                    "EKI 買掉的是頻率、留下的是嚴重度：機率降幅與預期損失降幅之間的楔子，就是客戶承接的尾部。")


def slide_monday(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, 8)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.5),
             [[("五  Monday 待決策", 24, True, BLACK)]])
    tasks = [
        "議題組合確認：1 主 + 2 補充（建議 01 存量健檢 + 05 防禦對照 + 06 TSLA 教材）",
        "重抓 yfinance 校正全部 est. 數字（本週 4 交易日已完整）",
        "存續 FCN 部位盤點：半導體單一股 vintage 清單 + 各筆距 KI 緩衝計算",
        "向交易台索取 NVDA/MU/TSLA/TSM 6M ATM IV 與 indicative coupon，替換示意值",
        "NLM saturation Round 1+ 加 Codex Gate 3 平行驗證",
        "寫 reading_map.md + 底稿.md（UWS 12 章節）",
        "build pptx/docx + lint_ryan_feedback 12 規則",
        "Codex audit + persona pre-screen（法人機構讀者視角）＋ quality_report A++ 簽核",
    ]
    y = Inches(1.45)
    for i, t in enumerate(tasks, 1):
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.42), Inches(0.42))
        circ.fill.solid()
        circ.fill.fore_color.rgb = GS_GOLD
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i)
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = BLACK
        r.font.name = EN_FONT
        add_text(slide, Inches(1.45), y + Inches(0.04), Inches(11.2), Inches(0.5),
                 [[(t, 13, False, DARK_GRAY)]])
        y += Inches(0.62)
    add_text(slide, Inches(0.7), Inches(6.9), Inches(12), Inches(0.35),
             [[(f"Generated {GEN_DATE} | 國泰證券法人債券業務部 · 法人業務二處", 10.5, False, MED_GRAY)]])


def slide_back(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLACK)
    add_rect(slide, Inches(0.7), Inches(3.3), Inches(2.2), Pt(3), GS_GOLD)
    add_text(slide, Inches(0.7), Inches(3.55), Inches(11.9), Inches(1.6),
             [[("FCN Weekly — KI/KO 期望值週報", 30, True, WHITE)],
              [(f"{ISSUE} 期  |  v0 partial scan  |  Monday 校正後定稿", 14, False, GS_GOLD)],
              [("本簡報之票息與條款均為示意，非報價；投資決策以正式產品文件為準。", 11, False, LIGHT_GRAY)]],
             line_spacing=1.5)


def build(out_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs)
    slide_framework(prs)
    slide_market(prs)
    slide_main_topics(prs)
    slide_supp_topics(prs)
    slide_matrix(prs)
    slide_eki_compare(prs)
    slide_monday(prs)
    slide_back(prs)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    out = Path(__file__).parent / "_draft_fcn_weekly_2026-07-03.pptx"
    build(str(out))
