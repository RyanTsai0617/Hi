"""ETF Weekly Friday topic preview — v2 PPTX skeleton with matplotlib charts + visual cards.

設計：9 slides、16:9、白底、teal #1B4F5A、微軟正黑體、footer 法人業務二處。
v2 升級（vs v1 文字表格）：
- Slide 2 反轉訊號改 matplotlib 雙 horizontal bar chart（替代 2 個表格）
- Slide 3-4 議題候選改 visual cards（彩色 box + headline + chips）
- Slide 5-6 ETF 矩陣加 heatmap 顏色強度 column
- Section divider 全頁過場（teal 大色塊 + yellow stripe）

本地測試用 0511 期 mock data、跑出 PPTX 看視覺、滿意後 update 進 routine prompt。
"""
from __future__ import annotations
import os
from pathlib import Path
from io import BytesIO

import matplotlib
matplotlib.use("Agg")  # no GUI
import matplotlib.pyplot as plt
from matplotlib import font_manager

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ============== Visual constants ==============
ACCENT = RGBColor(0x1B, 0x4F, 0x5A)         # TEAL_800
ACCENT_LIGHT = RGBColor(0x2B, 0x7A, 0x78)   # TEAL_600
BODY = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x4A, 0x4A, 0x4A)
META = RGBColor(0x76, 0x76, 0x76)
GAIN = RGBColor(0x15, 0x80, 0x3D)
LOSS = RGBColor(0xB9, 0x1C, 0x1C)
KPI_BG_1 = RGBColor(0xE8, 0xC5, 0x47)       # YELLOW
KPI_BG_2 = RGBColor(0x2B, 0x7A, 0x78)       # TEAL
KPI_BG_3 = RGBColor(0x1B, 0x4F, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
ALT_ROW = RGBColor(0xF2, 0xF4, 0xF5)
CARD_BG_1 = RGBColor(0xE8, 0xC5, 0x47)      # yellow
CARD_BG_2 = RGBColor(0x2B, 0x7A, 0x78)      # teal
CARD_BG_3 = RGBColor(0x1B, 0x4F, 0x5A)      # dark teal
CARD_BG_NEUTRAL = RGBColor(0xE5, 0xE5, 0xE5)
DIVIDER_YELLOW = RGBColor(0xE8, 0xC5, 0x47)

# Hex versions for matplotlib
ACCENT_HEX = "#1B4F5A"
ACCENT_LIGHT_HEX = "#2B7A78"
GAIN_HEX = "#15803D"
LOSS_HEX = "#B91C1C"
LIGHT_BG_HEX = "#F8F9FA"
META_HEX = "#767676"

ZH_FONT = "微軟正黑體"
EN_FONT = "Calibri"

# Try to use Microsoft JhengHei for matplotlib zh-rendering
for fpath in ["C:/Windows/Fonts/msjh.ttc", "C:/Windows/Fonts/msjh.ttf",
              "/Library/Fonts/Microsoft Sans Serif.ttf"]:
    if Path(fpath).exists():
        font_manager.fontManager.addfont(fpath)
        break
plt.rcParams["font.family"] = ["Microsoft JhengHei", "Microsoft YaHei", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False


# ============== python-pptx helpers ==============

def apply_font(run, *, size=14, bold=False, color=BODY, font_en=EN_FONT, font_zh=ZH_FONT):
    run.font.name = font_en
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_zh)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    paras = text.split("\n") if isinstance(text, str) else text
    for i, p_text in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = p_text
        apply_font(run, size=size, bold=bold, color=color)
    return tb


def add_rect(slide, x, y, w, h, fill, *, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
    return s


def add_rounded_rect(slide, x, y, w, h, fill):
    """Rounded rectangle for cards."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def add_image(slide, x, y, w, h, image_stream):
    """Add matplotlib figure as image."""
    return slide.shapes.add_picture(image_stream, x, y, width=w, height=h)


def add_footer(slide, num, total, footer_text):
    add_text(slide, Inches(0.6), Inches(7.1), Inches(11), Inches(0.3),
             footer_text, size=9, color=META,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(12.1), Inches(7.1), Inches(0.7), Inches(0.3),
             f"{num} / {total}", size=9, color=META,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_title_bar(slide, title, subtitle=None, x=Inches(0.6), y=Inches(0.4)):
    """Title + optional subtitle + accent underline."""
    add_text(slide, x, y, Inches(12), Inches(0.6), title,
             size=24, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, x, y + Inches(0.55), Inches(12), Inches(0.4), subtitle,
                 size=12, color=SUB)
    # Yellow accent underline
    add_rect(slide, x, y + Inches(1.0), Inches(0.6), Inches(0.05), DIVIDER_YELLOW)


# ============== Matplotlib chart helpers ==============

def make_reversal_bar_chart(rebound_rows, decline_rows, *, width_in=12.0, height_in=4.5):
    """Horizontal bar chart: top reversals up + top reversals down.

    rebound_rows: list of (etf, week_pct, prev_pct, reversal_pp) — positive reversal
    decline_rows: same, negative reversal
    Returns BytesIO PNG stream.
    """
    fig, axes = plt.subplots(1, 2, figsize=(width_in, height_in), dpi=150)
    fig.patch.set_facecolor("#FFFFFF")

    # Left: rebound (positive reversal, sorted highest at top)
    ax = axes[0]
    rebound = sorted(rebound_rows, key=lambda r: r[3])  # ascending so top is highest
    labels = [r[0] for r in rebound]
    values = [r[3] for r in rebound]
    bars = ax.barh(labels, values, color=GAIN_HEX, edgecolor="#0e5128", linewidth=0.5)
    for bar, v in zip(bars, values):
        ax.text(v + max(values) * 0.02, bar.get_y() + bar.get_height() / 2,
                f"+{v:.1f}", va="center", fontsize=9, color="#0e5128")
    ax.set_title("Top 反彈（本週強、上週弱）", fontsize=12, color=ACCENT_HEX, pad=8, weight="bold")
    ax.set_xlabel("反轉幅度 (pp)", fontsize=9, color=META_HEX)
    ax.set_facecolor(LIGHT_BG_HEX)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")
    ax.tick_params(colors=META_HEX)
    ax.grid(axis="x", linestyle=":", alpha=0.3)

    # Right: decline (negative reversal)
    ax = axes[1]
    decline = sorted(decline_rows, key=lambda r: r[3], reverse=True)  # descending so most negative at bottom (= top label)
    labels = [r[0] for r in decline]
    values = [r[3] for r in decline]
    bars = ax.barh(labels, values, color=LOSS_HEX, edgecolor="#7f1313", linewidth=0.5)
    for bar, v in zip(bars, values):
        ax.text(v - abs(min(values)) * 0.02, bar.get_y() + bar.get_height() / 2,
                f"{v:.1f}", va="center", ha="right", fontsize=9, color="#7f1313")
    ax.set_title("Top 反向下殺（本週弱、上週強）", fontsize=12, color=ACCENT_HEX, pad=8, weight="bold")
    ax.set_xlabel("反轉幅度 (pp)", fontsize=9, color=META_HEX)
    ax.set_facecolor(LIGHT_BG_HEX)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")
    ax.tick_params(colors=META_HEX)
    ax.grid(axis="x", linestyle=":", alpha=0.3)

    plt.tight_layout()
    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="#FFFFFF", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


def make_ytd_sparkline(ytd_values_dict, *, width_in=11.0, height_in=2.5):
    """Sparkline-style horizontal bar showing YTD% per ETF (top 12 in absolute magnitude)."""
    sorted_items = sorted(ytd_values_dict.items(), key=lambda kv: abs(kv[1]), reverse=True)[:12]
    sorted_items.reverse()  # so highest abs at top
    labels = [k for k, _ in sorted_items]
    values = [v for _, v in sorted_items]
    colors = [GAIN_HEX if v > 0 else LOSS_HEX for v in values]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=150)
    fig.patch.set_facecolor("#FFFFFF")
    bars = ax.barh(labels, values, color=colors, edgecolor="none")
    for bar, v in zip(bars, values):
        offset = max(abs(min(values)), abs(max(values))) * 0.02
        x_pos = v + offset if v > 0 else v - offset
        ha = "left" if v > 0 else "right"
        ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
                f"{'+' if v > 0 else ''}{v:.1f}%",
                va="center", ha=ha, fontsize=9,
                color=GAIN_HEX if v > 0 else LOSS_HEX)
    ax.set_title("YTD % — Top 12 by absolute magnitude", fontsize=11, color=ACCENT_HEX, pad=6, weight="bold")
    ax.set_facecolor(LIGHT_BG_HEX)
    ax.axvline(x=0, color="#666666", linewidth=0.6)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")
    ax.tick_params(colors=META_HEX, labelsize=9)
    ax.grid(axis="x", linestyle=":", alpha=0.3)

    plt.tight_layout()
    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="#FFFFFF", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


# ============== Slide builders ==============

def build_slide_cover(slide, period, week_range, status_text, kpi_cards, footer):
    # Background subtle: top teal stripe
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.4), ACCENT)
    add_rect(slide, Inches(0), Inches(0.4), Inches(13.333), Inches(0.05), DIVIDER_YELLOW)

    # Main title
    add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
             "美股 ETF 週報 — Friday topic preview",
             size=34, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Period
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.6),
             f"{period} 期 · 涵蓋 {week_range}",
             size=18, color=BODY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Status
    add_text(slide, Inches(0.6), Inches(3.0), Inches(12), Inches(0.5),
             status_text,
             size=12, color=SUB,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 3 KPI cards
    card_w = Inches(3.4)
    card_h = Inches(2.2)
    gap = Inches(0.3)
    cards_y = Inches(4.0)
    total_w = card_w * 3 + gap * 2
    start_x = (Inches(13.333) - total_w) / 2

    for i, (ticker, value, label, bg_color) in enumerate(kpi_cards):
        cx = start_x + i * (card_w + gap)
        # Card with rounded corners
        add_rounded_rect(slide, cx, cards_y, card_w, card_h, bg_color)
        # Ticker (top)
        add_text(slide, cx, cards_y + Inches(0.2), card_w, Inches(0.45), ticker,
                 size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Value (large, middle)
        add_text(slide, cx, cards_y + Inches(0.7), card_w, Inches(0.9), value,
                 size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Label (bottom)
        add_text(slide, cx, cards_y + Inches(1.6), card_w, Inches(0.5), label,
                 size=11, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 1, 9, footer)


def build_slide_section_divider(slide, section_num, section_title, footer):
    """Full-page teal background divider with large section number + title."""
    # Full bg
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), ACCENT)
    # Large section number (left, white, huge)
    add_text(slide, Inches(0.8), Inches(2.4), Inches(4), Inches(2.7),
             section_num, size=144, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Yellow stripe accent
    add_rect(slide, Inches(0.8), Inches(5.2), Inches(1.5), Inches(0.1), DIVIDER_YELLOW)
    # Section title (right of number)
    add_text(slide, Inches(5.5), Inches(3.0), Inches(7.0), Inches(1.5),
             section_title, size=36, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Footer (white tone since dark bg)
    add_text(slide, Inches(0.6), Inches(7.1), Inches(12.0), Inches(0.3),
             footer, size=9, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_slide_reversal_chart(slide, rebound_rows, decline_rows, footer):
    """Slide 2 (or 3): horizontal bar chart of reversals."""
    add_title_bar(slide, "反轉訊號", "本週 partial vs 上週、Top 10 反彈（綠）+ Top 10 反向下殺（紅）")

    chart_buf = make_reversal_bar_chart(rebound_rows, decline_rows,
                                          width_in=12.0, height_in=4.7)
    add_image(slide, Inches(0.6), Inches(1.6), Inches(12.0), Inches(4.7), chart_buf)

    # Insight callout
    top_rebound = max(rebound_rows, key=lambda r: r[3])
    top_decline = min(decline_rows, key=lambda r: r[3])
    callout = (f"最強反彈：{top_rebound[0]} {top_rebound[3]:+.1f} pp（本週 {top_rebound[1]:+.2f}% / 上週 {top_rebound[2]:+.2f}%）"
               f"　|　最猛下殺：{top_decline[0]} {top_decline[3]:+.1f} pp（本週 {top_decline[1]:+.2f}% / 上週 {top_decline[2]:+.2f}%）")
    add_text(slide, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.5),
             callout, size=11, color=ACCENT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 3, 9, footer)


def build_slide_topic_cards(slide, title, subtitle, candidates, *, page_num, footer, card_bg_colors):
    """議題候選 visual cards (3 cards on slide).

    candidates: list of dicts with keys: 主標, 核心數據, 相關 ETF, driver, 對台灣機構意義
    """
    add_title_bar(slide, title, subtitle)

    # 3 cards stacked vertically
    n = len(candidates)
    card_h_total = Inches(5.0)
    card_h = card_h_total / n - Inches(0.15)
    card_y_start = Inches(1.6)
    card_w = Inches(12.0)
    card_x = Inches(0.6)

    for i, c in enumerate(candidates):
        cy = card_y_start + i * (card_h + Inches(0.15))
        bg = card_bg_colors[i % len(card_bg_colors)]
        # Card bg (light tint behind)
        add_rounded_rect(slide, card_x, cy, card_w, card_h, LIGHT_BG)
        # Left accent stripe (color)
        add_rect(slide, card_x, cy, Inches(0.2), card_h, bg)
        # Headline
        add_text(slide, card_x + Inches(0.4), cy + Inches(0.1), card_w - Inches(0.6), Inches(0.4),
                 c.get("主標", ""), size=14, bold=True, color=ACCENT)
        # Metadata row 1: 核心數據 + 相關 ETF
        meta1 = f"核心：{c.get('核心數據', '')}　·　相關 ETF：{c.get('相關 ETF', '')}"
        add_text(slide, card_x + Inches(0.4), cy + Inches(0.55), card_w - Inches(0.6), Inches(0.35),
                 meta1, size=10, color=BODY)
        # Metadata row 2: driver + 對台灣意義
        meta2 = f"driver：{c.get('driver', '')}　·　對台灣機構意義：{c.get('對台灣機構意義', '')}"
        add_text(slide, card_x + Inches(0.4), cy + Inches(0.95), card_w - Inches(0.6), Inches(0.6),
                 meta2, size=9.5, color=SUB)

    add_footer(slide, page_num, 9, footer)


def build_slide_etf_matrix(slide, theme_name, theme_rows, *, page_num, footer):
    """ETF 矩陣 with heatmap-style coloring on week% column.

    theme_rows: list of (代號, 全名, 一週%, 上週%, YTD%, 角色定位)
    """
    add_title_bar(slide, f"ETF 矩陣 — {theme_name}", "一週% 含正負色標、YTD% 含絕對值色標")

    # Build table
    headers = ["代號", "全名", "一週 %", "上週 %", "YTD %", "角色定位"]
    n_cols = len(headers)
    n_rows = 1 + len(theme_rows)
    tbl_x = Inches(0.4)
    tbl_y = Inches(1.6)
    tbl_w = Inches(12.5)
    tbl_h = Inches(5.2)

    table = slide.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y, tbl_w, tbl_h).table

    # Header
    for j, hd in enumerate(headers):
        c = table.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
        c.text_frame.clear()
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hd
        apply_font(r, size=11, bold=True, color=WHITE)

    # Body
    import re
    for i, row in enumerate(theme_rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            t = str(val).strip()
            # Heatmap on 一週% (col 2) and YTD% (col 4)
            if j in (2, 4) and re.match(r"^[\-+]?\d", t):
                pct_val = float(t.replace("%", "").replace("+", ""))
                # Color intensity by magnitude
                if pct_val > 0:
                    intensity = min(abs(pct_val) / 8.0, 1.0)  # cap at 8% = full intensity
                    g = int(0xC0 - 0x60 * intensity)
                    bg = RGBColor(0xE8 - int(0x40 * intensity), 0xF5 - int(0x30 * intensity), 0xE8 - int(0x40 * intensity))
                    c.fill.solid()
                    c.fill.fore_color.rgb = bg
                elif pct_val < 0:
                    intensity = min(abs(pct_val) / 8.0, 1.0)
                    bg = RGBColor(0xFD - int(0x10 * intensity), 0xE5 - int(0x40 * intensity), 0xE5 - int(0x40 * intensity))
                    c.fill.solid()
                    c.fill.fore_color.rgb = bg
            elif i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = ALT_ROW

            c.text_frame.clear()
            p = c.text_frame.paragraphs[0]
            if re.match(r"^[\-+\$]?\d", t) or t.endswith("%") or t.lower() in ("n/a", "na"):
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = t
            color = BODY
            if t.startswith("+") and re.match(r"^\+\d", t):
                color = GAIN
            elif t.startswith("-") and re.match(r"^-\d", t):
                color = LOSS
            apply_font(r, size=9.5, color=color)

    add_footer(slide, page_num, 9, footer)


def build_slide_macro_research(slide, macro_events, research_reports, *, page_num, footer):
    add_title_bar(slide, "Macro 事件 + A 級 research", "本週主要 catalyst")

    # Left col: macro events
    add_text(slide, Inches(0.6), Inches(1.6), Inches(6.0), Inches(0.4),
             "📌 Macro 事件本週", size=14, bold=True, color=ACCENT_LIGHT)
    macro_text = "\n".join([f"•  {e}" for e in macro_events])
    add_text(slide, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.5),
             macro_text, size=11, color=BODY)

    # Right col: A 級 research
    add_text(slide, Inches(7.0), Inches(1.6), Inches(5.7), Inches(0.4),
             "📚 A 級 research（本週）", size=14, bold=True, color=ACCENT_LIGHT)
    if research_reports:
        research_text = "\n".join([f"•  {r}" for r in research_reports])
    else:
        research_text = "(本週無重要 A 級 research、或 yfinance scan 未觸發)"
    add_text(slide, Inches(7.0), Inches(2.1), Inches(5.7), Inches(4.5),
             research_text, size=11, color=BODY)

    add_footer(slide, page_num, 9, footer)


def build_slide_monday_tasks(slide, period, tasks, footer):
    add_title_bar(slide, f"Monday {period} 待決策", "週一早上接續、預估 2-2.5 hr 全套完成")

    # Tasks as numbered cards
    n = len(tasks)
    task_h = Inches(0.6)
    task_y_start = Inches(1.7)
    task_w = Inches(12.0)
    task_x = Inches(0.6)

    for i, task in enumerate(tasks):
        ty = task_y_start + i * (task_h + Inches(0.05))
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, task_x, ty + Inches(0.05),
                                         Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        # Number text in circle
        add_text(slide, task_x, ty + Inches(0.05), Inches(0.5), Inches(0.5),
                 str(i + 1), size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Task text
        add_text(slide, task_x + Inches(0.7), ty + Inches(0.05), task_w - Inches(0.7), Inches(0.5),
                 task, size=12, color=BODY, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 9, 9, footer)


# ============== Main builder ==============

def build_friday_preview(*,
                          this_friday: str,
                          period: str,
                          week_range: str,
                          status_text: str,
                          kpi_cards: list,
                          rebound_rows: list,
                          decline_rows: list,
                          main_candidates: list,
                          sub_candidates: list,
                          theme_tables: list,
                          macro_events: list,
                          research_reports: list,
                          monday_tasks: list,
                          out_path: str):
    """Build the v2 Friday preview PPTX.

    All inputs explicit for testability.
    """
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    blank = pres.slide_layouts[6]
    footer = f"{this_friday} · 國泰證券法人債券業務部 · 法人業務二處 · Friday topic preview"

    # 1. Cover
    s1 = pres.slides.add_slide(blank)
    build_slide_cover(s1, period, week_range, status_text, kpi_cards, footer)

    # 2. Section divider 01: Reversal
    s2 = pres.slides.add_slide(blank)
    build_slide_section_divider(s2, "01", "反轉訊號 · Reversal", footer)

    # 3. Reversal chart
    s3 = pres.slides.add_slide(blank)
    build_slide_reversal_chart(s3, rebound_rows, decline_rows, footer)

    # 4. 主議題候選 cards
    s4 = pres.slides.add_slide(blank)
    build_slide_topic_cards(s4, "主議題候選 · Top 3", "敘事最強、優先選為主議題",
                             main_candidates, page_num=4, footer=footer,
                             card_bg_colors=[CARD_BG_3, CARD_BG_2, CARD_BG_1])

    # 5. 補充議題候選 cards
    s5 = pres.slides.add_slide(blank)
    build_slide_topic_cards(s5, "補充議題候選 · Top 5", "補充視角、選 2 個搭配主議題",
                             sub_candidates, page_num=5, footer=footer,
                             card_bg_colors=[CARD_BG_2, CARD_BG_1, CARD_BG_3, CARD_BG_2, CARD_BG_1])

    # 6-7. ETF 矩陣 (1-2 themes)
    page_num = 6
    for theme_name, theme_rows in theme_tables[:2]:
        s = pres.slides.add_slide(blank)
        build_slide_etf_matrix(s, theme_name, theme_rows, page_num=page_num, footer=footer)
        page_num += 1

    # 8. Macro + research
    s8 = pres.slides.add_slide(blank)
    build_slide_macro_research(s8, macro_events, research_reports, page_num=8, footer=footer)

    # 9. Monday tasks
    s9 = pres.slides.add_slide(blank)
    build_slide_monday_tasks(s9, period, monday_tasks, footer)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    pres.save(out_path)
    print(f"[OK] Friday preview v2 saved: {out_path}")
    return out_path


# ============== Mock data for local test (0511 期) ==============

if __name__ == "__main__":
    # Mock 0511 期 data based on what we computed manually 5/8 earlier
    out_path = r"C:/Users/tsait/OneDrive/桌面/每周報告/20260511/_v2_friday_preview_test.pptx"

    build_friday_preview(
        this_friday="2026-05-08",
        period="20260511",
        week_range="2026-05-04 ~ 2026-05-08",
        status_text="v0 partial scan · Mon-Thu close (4 trading days) · Monday 重抓含週五 close",
        kpi_cards=[
            ("IRBO", "+7.99%", "本週最強 · AI Brain 鍊條", CARD_BG_1),
            ("USO", "-5.48%", "YTD 跌穿百分百 · 油市急殺", CARD_BG_2),
            ("GDX", "+5.29%", "V 字反彈 · 黃金回神", CARD_BG_3),
        ],
        rebound_rows=[
            ("GDX", 5.29, -7.66, 12.95),
            ("SLV", 4.85, -0.73, 5.58),
            ("SMH", 5.94, 0.67, 5.27),
            ("SOXX", 5.71, 0.90, 4.81),
            ("XBI", 2.48, -2.15, 4.63),
            ("GLD", 2.01, -2.32, 4.33),
            ("XLK", 4.83, 1.03, 3.80),
            ("EEM", 3.84, 0.61, 3.23),
            ("IBB", 1.69, -1.48, 3.17),
            ("ITA", 2.89, 0.22, 2.67),
        ],
        decline_rows=[
            ("BNO", -6.29, 8.90, -15.19),
            ("UCO", -6.18, 8.97, -15.15),
            ("USO", -5.48, 7.85, -13.33),
            ("XOP", -6.01, 5.13, -11.14),
            ("XLE", -4.93, 3.48, -8.41),
            ("ICLN", -1.15, 4.80, -5.95),
            ("DBC", -1.82, 3.18, -5.00),
            ("UNG", -0.28, 3.88, -4.16),
            ("XLU", -3.07, 0.80, -3.87),
            ("AMLP", -0.48, 3.31, -3.79),
        ],
        main_candidates=[
            {
                "主標": "B+G AI/Robotics 鍊條深化 — 半導體 reignite + 人形機器人浮現",
                "核心數據": "SMH +5.94% · IRBO +7.99% · SOXX +5.71%",
                "相關 ETF": "SMH, SOXX, IRBO, THNQ, AIQ, XLK, BOTZ, ROBO, KWEB",
                "driver": "AMD Q1 +38% / TSMC 52w 高 / Morgan Stanley 5/7 humanoid 報告",
                "對台灣機構意義": "TSMC / HIWIN 在 Humanoid 100 名單、5/22 NVDA 為下一驗證",
            },
            {
                "主標": "A 油市急殺 — Hormuz 重啟談判 + Iran 提案 trigger",
                "核心數據": "USO -5.48% (YTD 跌穿百分百) · BNO -6.29% · XOP -6.01%",
                "相關 ETF": "USO, BNO, UCO, XOP, XLE, AMLP",
                "driver": "Trump 提案 + Iran ceasefire + Brent $115→$103.54 (-5.8%)",
                "對台灣機構意義": "油價回落減輕通膨壓力、Fed 降息預期可能回升",
            },
            {
                "主標": "D 黃金 V 字反彈 — 油跌帶實質殖利率回落",
                "核心數據": "GDX +5.29% (前週 -7.66%) · GLD +2.01% · SLV +4.85%",
                "相關 ETF": "GDX, GLD, SLV",
                "driver": "油跌 → 通膨預期降 → 實質殖利率回落 → 黃金 carry 優勢",
                "對台灣機構意義": "0504 期黃金崩盤完美反向、商品內部分化主軸不變",
            },
        ],
        sub_candidates=[
            {
                "主標": "新興市場補漲 — EEM/KWEB 美元疲弱受惠",
                "核心數據": "EEM +3.84% · KWEB +2.57% · FXI +1.03%",
                "相關 ETF": "EEM, KWEB, MCHI, FXI, INDA",
                "driver": "美元疲軟 + 油價跌 + MS 報告強調中國 humanoid VC 占 46%",
                "對台灣機構意義": "新興市場 ETF 配置時點可能成熟",
            },
            {
                "主標": "Biotech / 醫療補漲 — XBI/IBB 落後類股回神",
                "核心數據": "XBI +2.48% (前週 -2.15%) · IBB +1.69% (前週 -1.48%)",
                "相關 ETF": "XBI, IBB, XLV",
                "driver": "落後類股輪動、AI / 半導體 spill-over",
                "對台灣機構意義": "防禦型生技配置時機觀察",
            },
            {
                "主標": "Stagflation 訊號 — ISM Prices 70.7 黏住",
                "核心數據": "ISM Services 53.6 / New Orders -7.1 / Prices 70.7",
                "相關 ETF": "TLT, IEF, TIP, XLU, XLP",
                "driver": "Fed 6 月降息機率 5.1% (CME)、stagflation confirmed",
                "對台灣機構意義": "TLT 配置繼續觀望、防禦類股 XLU 反弱訊號",
            },
            {
                "主標": "國防補漲 — ITA + ICE 地緣議題延續",
                "核心數據": "ITA +2.89% (前週 +0.22%)",
                "相關 ETF": "ITA, XAR",
                "driver": "Hormuz 雖緩解但長線地緣議題仍在",
                "對台灣機構意義": "國防 ETF 加碼考量",
            },
            {
                "主標": "公用 / 防禦反弱 — XLU/XLP 與市場分化",
                "核心數據": "XLU -3.07% (前週 +0.80%) · XLP / XLV 弱",
                "相關 ETF": "XLU, XLP, XLV, RSP",
                "driver": "Narrow leadership 集中半導體 + AI、其他類股弱",
                "對台灣機構意義": "等權 RSP 可能跑輸 SPY、配置需檢視 cap-weight bias",
            },
        ],
        theme_tables=[
            ("AI / 半導體 / Robotics", [
                ["SMH", "VanEck 半導體 (含 TSM)", "+5.94%", "+0.67%", "+48.29%", "0511 主議題核心"],
                ["SOXX", "iShares 半導體純美股", "+5.71%", "+0.90%", "+61.36%", "YTD 創新高"],
                ["IRBO", "iShares Robotics & AI", "+7.99%", "+1.19%", "+39.53%", "本週最強"],
                ["THNQ", "ROBO Global AI", "+6.65%", "+1.77%", "+22.26%", "AI 純"],
                ["AIQ", "Global X AI", "+5.58%", "+3.05%", "+16.61%", "AI 純廣"],
                ["XLK", "科技類股", "+4.83%", "+1.03%", "+16.47%", "大型科技"],
                ["BOTZ", "Global X Robotics", "+3.19%", "+3.32%", "+8.85%", "工業機器人"],
                ["ROBO", "ROBO Global Robotics", "+2.28%", "+2.44%", "+20.45%", "全 Robotics"],
                ["KWEB", "中國互聯網", "+2.57%", "-0.17%", "-14.36%", "中國 humanoid 投資鍊"],
            ]),
            ("商品 / 避險 / 大盤", [
                ["USO", "WTI 期貨", "-5.48%", "+7.85%", "+93.89%", "0504 主議題反向"],
                ["BNO", "Brent 期貨", "-6.29%", "+8.90%", "+88.51%", "對沖中東反向"],
                ["XOP", "油氣 E&P", "-6.01%", "+5.13%", "+31.84%", "高 Beta"],
                ["XLE", "綜合能源", "-4.93%", "+3.48%", "+26.21%", "長線部位仍佳"],
                ["GDX", "黃金礦業", "+5.29%", "-7.66%", "+6.84%", "V 字反彈"],
                ["GLD", "黃金實體", "+2.01%", "-2.32%", "+8.30%", "實質殖利率回落"],
                ["SLV", "白銀", "+4.85%", "-0.73%", "+8.47%", "工業屬性 + 避險"],
                ["TLT", "20+ 年公債", "+0.05%", "-0.90%", "-1.30%", "Fed 預期 stable"],
                ["SPY", "標普 500", "+1.52%", "+0.94%", "+6.65%", "5/7 創 YTD 新高"],
            ]),
        ],
        macro_events=[
            "Hormuz 重啟談判 (Trump 提案 + Iran ceasefire 維持) — Brent $115→$103.54 / -5.8%",
            "AMD Q1 2026 營收 $10.3B (+38% YoY) — ignited semiconductor rally",
            "TSMC 5/7 +6% 至 $418 — 52-week high",
            "ISM Services PMI 4 月 53.6 (連 22 月擴張) — Prices 70.7 黏住、stagflation confirmed",
            "Fed 6 月降息機率 5.1% (CME FedWatch 5/7) — 94.9% 預期 hold 3.50-3.75%",
            "5/8 NFP (ET 8:30am 公布) — Monday patch 用",
            "Gold spot 5/7 $4,705 — 自 4/28 $4,620 反彈 +1.8%",
        ],
        research_reports=[
            "Morgan Stanley 5/7 「Humanoids: Humanoid Horizons: Money Meets Machines」 — Tech Diffusion 2026 Key Theme",
            "  · 2026 YTD VC 融資已超 2025 全年 (China 占 46%、April 41 deals)",
            "  · Bezos Project Prometheus $38B / Meta x ARI / Unitree IPO",
            "  · Humanoid 100 Index 自 2025/2/6 +45% (跑贏 SPX/MSCI)",
        ],
        monday_tasks=[
            "議題組合確認：1 主 + 2 補充（從 8 候選圈選）",
            "重抓 yfinance 含 5/8 close — 5 trading days 完整週",
            "NLM saturation Round 1+ 加 Codex Gate 3 平行驗證",
            "寫 reading_map.md + 底稿.md（UWS 12 章節）",
            "build docx + lint_ryan_feedback 12 規則",
            "NLM Studio Path A 全來源 + Path B 只 docx + Codex Path C deterministic",
            "Codex audit Path A/B/C 一致性 + Codex Ryan persona pre-screen",
            "quality_report A++ 簽核 + reviewed.docx canonical",
        ],
        out_path=out_path,
    )
