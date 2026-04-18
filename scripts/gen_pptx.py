#!/usr/bin/env python3
"""
PDF-to-Briefing PPTX Generator — v3 (Weekly20260406 Style)
Matches the design system of Weekly20260406.pptx:
  - 4:3, teal #1B7A6E top bar, 48pt title, subtitle bar, card layout
  - Color-coded category bars, white card containers, footer + page number

Usage:
    python gen_pptx.py --input slides.json --output briefing.pptx
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette (from Weekly20260406.pptx) ───────────────────────
TEAL       = RGBColor(0x1B, 0x7A, 0x6E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x40, 0x40, 0x40)
SUBTITLE_BG = RGBColor(0x50, 0x50, 0x50)
GRAY_FOOT  = RGBColor(0x77, 0x77, 0x77)
RED        = RGBColor(0xC0, 0x39, 0x2B)
BLUE       = RGBColor(0x29, 0x80, 0xB9)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
GREEN      = RGBColor(0x00, 0xA0, 0x00)
CARD_BG    = RGBColor(0xF8, 0xF9, 0xFA)

# Category colors for section slides (rotate through)
SECTION_COLORS = [TEAL, RED, BLUE, ORANGE]
_section_idx = 0

# ── Font ───────────────────────────────────────────────────────────
FONT = "Microsoft JhengHei"

# ── Slide dimensions (4:3) ─────────────────────────────────────────
SLIDE_W = 9144000
SLIDE_H = 6858000

# ── Layout constants (EMU) ─────────────────────────────────────────
TOP_BAR_H    = 54864       # thin teal top bar
TITLE_LEFT   = 628650
TITLE_TOP    = 320674
TITLE_W      = 7886700
TITLE_H      = 1325563
SUBTITLE_TOP = 1323071
SUBTITLE_H   = 461665
BODY_TOP     = 2066544     # below subtitle bar
BODY_H       = SLIDE_H - BODY_TOP - 500000
FOOTER_TOP   = 6492240
FOOTER_COMP  = "國泰證券｜法人金融事業處"
DIVIDER_LEFT = 3200400
DIVIDER_W    = 2743200

_page_num = 0


def _next_page():
    global _page_num
    _page_num += 1
    return _page_num


def _next_section_color():
    global _section_idx
    c = SECTION_COLORS[_section_idx % len(SECTION_COLORS)]
    _section_idx += 1
    return c


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Helpers ────────────────────────────────────────────────────────

def add_top_bar(slide):
    """Thin teal line at the very top."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TOP_BAR_H)
    s.fill.solid(); s.fill.fore_color.rgb = TEAL
    s.line.fill.background()


def add_title(slide, text):
    """48pt teal bold title."""
    tb = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT; p.font.size = Pt(48); p.font.bold = True
    p.font.color.rgb = TEAL


def add_subtitle_bar(slide, text, color=SUBTITLE_BG):
    """Subtitle text bar below title."""
    tb = slide.shapes.add_textbox(0, SUBTITLE_TOP, SLIDE_W, SUBTITLE_H)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT; p.font.size = Pt(24); p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    # indent
    p.font.size = Pt(24)


def add_divider(slide):
    """Centered teal divider line below title."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, DIVIDER_LEFT, 749808, DIVIDER_W, 18288)
    s.fill.solid(); s.fill.fore_color.rgb = TEAL
    s.line.fill.background()


def add_footer(slide, page_no, source="資料來源：各券商研究報告"):
    """Footer with source and page number."""
    # Source
    tb = slide.shapes.add_textbox(274320, FOOTER_TOP, 6400800, 228600)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = source
    p.font.name = FONT; p.font.size = Pt(9); p.font.color.rgb = GRAY_FOOT

    # Page number
    tb2 = slide.shapes.add_textbox(7772400, FOOTER_TOP, 1097280, 228600)
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_no)
    p2.font.name = FONT; p2.font.size = Pt(11); p2.font.color.rgb = TEAL
    p2.alignment = PP_ALIGN.RIGHT


def _textbox(slide, left, top, width, height, text,
             size=Pt(16), color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT; p.font.size = size
    p.font.color.rgb = color; p.font.bold = bold
    p.alignment = align
    return tb


def _bullet_list(slide, left, top, width, height, bullets,
                 size=Pt(16), color=DARK_TEXT, spacing=Pt(8)):
    # Auto-shrink: if too many bullets, reduce font to avoid overflow
    n = len(bullets)
    if n > 8 and size >= Pt(16):
        size = Pt(13)
        spacing = Pt(4)
    elif n > 12:
        size = Pt(11)
        spacing = Pt(3)

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {b}"
        p.font.name = FONT; p.font.size = size; p.font.color.rgb = color
        p.space_after = spacing
    return tb


def _card_box(slide, left, top, width, height, fill_color=WHITE, border_color=None):
    """White card container with subtle border."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


# ── Slide builders ─────────────────────────────────────────────────

def build_title_slide(prs, data):
    """Cover slide using theme layout (simulated)."""
    slide = _blank(prs)

    # Title
    _textbox(slide, 685800, 1122363, 4384964, 2306637,
             data.get("title", "Briefing"),
             size=Pt(48), color=TEAL, bold=True)

    # Subtitle
    sub = data.get("subtitle", "")
    if sub:
        _textbox(slide, 685799, 3518910, 4384963, 1161155,
                 sub, size=Pt(24), color=SUBTITLE_BG, bold=True)

    # Focus highlights (e.g. "本週焦點：週四 CPI、週五 核心 PCE")
    highlights = data.get("highlights", [])
    if highlights:
        hl_text = "  |  ".join(highlights) if isinstance(highlights, list) else str(highlights)
        tb = _textbox(slide, 685800, 4800000, 7772400, 457200,
                      hl_text, size=Pt(14), color=RED, bold=True)
        # Accent line above highlights
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 685800, 4700000, 2743200, 18288)
        s.fill.solid(); s.fill.fore_color.rgb = TEAL
        s.line.fill.background()


def build_section_slide(prs, slide_data):
    """Section divider - full teal background with white title."""
    slide = _blank(prs)
    # Full teal bg
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = TEAL
    s.line.fill.background()

    _textbox(slide, TITLE_LEFT, Emu(2500000), TITLE_W, TITLE_H,
             slide_data.get("title", ""),
             size=Pt(44), color=WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Bottom accent line
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TITLE_LEFT, Emu(3700000), Emu(2000000), Emu(18288))
    s2.fill.solid(); s2.fill.fore_color.rgb = WHITE
    s2.line.fill.background()


def build_content_slide(prs, slide_data):
    """Standard content: top bar + title + subtitle bar + bullet list."""
    slide = _blank(prs)
    add_top_bar(slide)
    pn = _next_page()

    add_title(slide, slide_data.get("title", ""))
    add_divider(slide)

    # Subtitle if present
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_subtitle_bar(slide, subtitle, SUBTITLE_BG)

    # Bullets in a white card
    bullets = slide_data.get("bullets", [])
    if bullets:
        card_top = SUBTITLE_TOP + SUBTITLE_H + 91440 if subtitle else BODY_TOP - 200000
        card_h = FOOTER_TOP - card_top - 182880
        _card_box(slide, 365760, card_top, 8412480, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
        _bullet_list(slide, 548640, card_top + 137160, 8046720, card_h - 274320,
                     bullets, size=Pt(16), color=DARK_TEXT)

    add_footer(slide, pn)


def build_two_column_slide(prs, slide_data):
    """Two-column with card containers."""
    slide = _blank(prs)
    add_top_bar(slide)
    pn = _next_page()

    add_title(slide, slide_data.get("title", ""))
    add_divider(slide)

    card_top = BODY_TOP - 200000
    card_h = FOOTER_TOP - card_top - 182880
    col_w = 3977640
    gap = 274320

    # Left card
    _card_box(slide, 457200, card_top, col_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    lt = slide_data.get("left_title", "")
    if lt:
        _textbox(slide, 566928, card_top + 54864, Emu(col_w - 200000), 365760,
                 lt, size=Pt(18), color=TEAL, bold=True)
    lb = slide_data.get("left_bullets", [])
    if lb:
        _bullet_list(slide, 566928, card_top + 457200, Emu(col_w - 200000), Emu(card_h - 600000),
                     lb, size=Pt(13), color=DARK_TEXT)

    # Right card
    right_x = 457200 + col_w + gap
    _card_box(slide, right_x, card_top, col_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    rt = slide_data.get("right_title", "")
    if rt:
        _textbox(slide, right_x + 109728, card_top + 54864, Emu(col_w - 200000), 365760,
                 rt, size=Pt(18), color=RED, bold=True)
    rb = slide_data.get("right_bullets", [])
    if rb:
        _bullet_list(slide, right_x + 109728, card_top + 457200, Emu(col_w - 200000), Emu(card_h - 600000),
                     rb, size=Pt(13), color=DARK_TEXT)

    add_footer(slide, pn)


def build_discussion_slide(prs, slide_data):
    """Discussion questions in card layout."""
    slide = _blank(prs)
    add_top_bar(slide)
    pn = _next_page()

    add_title(slide, slide_data.get("title", "討論議題"))
    add_divider(slide)

    questions = slide_data.get("questions", [])
    y = BODY_TOP - 100000
    card_w = 8229600

    for i, q in enumerate(questions):
        question = q.get("question", "")
        context = q.get("context", "")

        card_h = 640080 if context else 365760
        _card_box(slide, 457200, y, card_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))

        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 502920, y + 54864, 274320, 274320)
        badge.fill.solid(); badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        _textbox(slide, 502920, y + 54864, 274320, 274320,
                 str(i + 1), size=Pt(14), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Question
        _textbox(slide, 868680, y + 54864, Emu(card_w - 600000), 274320,
                 question, size=Pt(16), color=TEAL, bold=True)

        if context:
            _textbox(slide, 868680, y + 365760, Emu(card_w - 600000), 228600,
                     context, size=Pt(11), color=GRAY_FOOT)

        y += card_h + 91440

    add_footer(slide, pn)


def build_image_slide(prs, slide_data):
    """Title + subtitle bar + full-width image (matches Weekly style)."""
    slide = _blank(prs)
    add_top_bar(slide)
    pn = _next_page()

    add_title(slide, slide_data.get("title", ""))

    # Subtitle/caption as subtitle bar
    caption = slide_data.get("caption", "")
    if caption:
        add_subtitle_bar(slide, caption, SUBTITLE_BG)

    image_path = slide_data.get("image_path", "")
    if image_path and Path(image_path).exists():
        img_top = BODY_TOP
        img_h = FOOTER_TOP - img_top - 182880
        slide.shapes.add_picture(image_path, 0, img_top, SLIDE_W, img_h)
    else:
        _textbox(slide, Emu(2000000), Emu(3500000), Emu(5000000), Emu(400000),
                 f"[Image not found: {image_path}]",
                 size=Pt(12), color=GRAY_FOOT, align=PP_ALIGN.CENTER)

    add_footer(slide, pn)


def build_image_with_text_slide(prs, slide_data):
    """Image on left + bullets on right in cards."""
    slide = _blank(prs)
    add_top_bar(slide)
    pn = _next_page()

    add_title(slide, slide_data.get("title", ""))
    add_divider(slide)

    image_path = slide_data.get("image_path", "")
    bullets = slide_data.get("bullets", [])
    card_top = BODY_TOP - 200000
    card_h = FOOTER_TOP - card_top - 182880

    if image_path and Path(image_path).exists():
        from PIL import Image as PILImage
        with PILImage.open(image_path) as img:
            iw, ih = img.size

        max_w = Emu(4200000)
        max_h = Emu(card_h)
        scale = min(max_w / Emu(int(iw * 914400 / 96)),
                    max_h / Emu(int(ih * 914400 / 96)))
        if scale > 1: scale = 1

        pw = int(iw * 914400 / 96 * scale)
        ph = int(ih * 914400 / 96 * scale)
        slide.shapes.add_picture(image_path, 365760, card_top, pw, ph)

        if bullets:
            right_x = 365760 + pw + 274320
            right_w = SLIDE_W - right_x - 365760
            _card_box(slide, right_x, card_top, right_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
            _bullet_list(slide, right_x + 137160, card_top + 137160,
                         right_w - 274320, card_h - 274320,
                         bullets, size=Pt(13))
    else:
        if bullets:
            _card_box(slide, 365760, card_top, 8412480, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
            _bullet_list(slide, 548640, card_top + 137160, 8046720, card_h - 274320,
                         bullets, size=Pt(16))

    add_footer(slide, pn)


def build_summary_slide(prs, slide_data):
    """Key takeaway / conclusion slide — teal accent with numbered items."""
    slide = _blank(prs)
    add_top_bar(slide)
    pn = _next_page()

    add_title(slide, slide_data.get("title", "本週結論"))
    add_divider(slide)

    items = slide_data.get("items", slide_data.get("bullets", []))
    y = BODY_TOP - 100000
    card_w = 8229600

    for i, item in enumerate(items):
        text = item if isinstance(item, str) else item.get("text", "")
        label = item.get("label", "") if isinstance(item, dict) else ""

        card_h = Emu(548640)
        _card_box(slide, 457200, y, card_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))

        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 502920, y + 54864, 274320, 274320)
        badge.fill.solid(); badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        _textbox(slide, 502920, y + 54864, 274320, 274320,
                 str(i + 1), size=Pt(16), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Label (e.g. "操作建議") if present
        if label:
            _textbox(slide, 868680, y + 36576, Emu(card_w - 600000), 228600,
                     label, size=Pt(11), color=TEAL, bold=True)
            _textbox(slide, 868680, y + 265176, Emu(card_w - 600000), 274320,
                     text, size=Pt(14), color=DARK_TEXT)
        else:
            _textbox(slide, 868680, y + 91440, Emu(card_w - 600000), Emu(card_h - 182880),
                     text, size=Pt(15), color=DARK_TEXT)

        y += card_h + 73152

    add_footer(slide, pn)


def build_end_slide(prs):
    """End slide - teal background."""
    slide = _blank(prs)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = TEAL
    s.line.fill.background()

    _textbox(slide, Emu(685800), Emu(2200000), Emu(7772400), Emu(1097280),
             "Thank You",
             size=Pt(48), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _textbox(slide, Emu(685800), Emu(3400000), Emu(7772400), Emu(640080),
             FOOTER_COMP,
             size=Pt(20), color=RGBColor(0xB2, 0xDF, 0xDB), align=PP_ALIGN.CENTER)

    _textbox(slide, Emu(685800), Emu(4000000), Emu(7772400), Emu(457200),
             "Claude + NotebookLM Co-Reading Analysis",
             size=Pt(14), color=RGBColor(0xB2, 0xDF, 0xDB), align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────

SLIDE_BUILDERS = {
    "section":         build_section_slide,
    "content":         build_content_slide,
    "two_column":      build_two_column_slide,
    "discussion":      build_discussion_slide,
    "image":           build_image_slide,
    "image_with_text": build_image_with_text_slide,
    "summary":         build_summary_slide,
}


def generate(input_path, output_path):
    global _page_num, _section_idx
    _page_num = 0
    _section_idx = 0

    with open(input_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = create_presentation()
    build_title_slide(prs, data)

    for sd in data.get("slides", []):
        st = sd.get("type", "content")
        builder = SLIDE_BUILDERS.get(st, build_content_slide)
        builder(prs, sd)

    build_end_slide(prs)

    prs.save(output_path)
    print(f"PPTX saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


def main():
    parser = argparse.ArgumentParser(description="Generate briefing PPTX (Weekly v3 style)")
    parser.add_argument("--input", required=True, help="Input JSON file path")
    parser.add_argument("--output", required=True, help="Output PPTX file path")
    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f"Error: Input file not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    generate(args.input, args.output)


if __name__ == "__main__":
    main()
