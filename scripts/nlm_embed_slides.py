#!/usr/bin/env python3
"""
NLM Embed Slides — 裁剪 NLM 投影片並嵌入 Weekly PPTX

流程：
  1. 讀取 NLM 生成的投影片 PDF
  2. 每頁動態偵測標題列 + 浮水印位置
  3. 裁剪（去掉上方標題列 + 下方 NotebookLM 浮水印）
  4. 用底稿 PPTX 的 Layout 5「只有標題」嵌入，副標題 = NLM 頁面標題

用法：
  python3 nlm_embed_slides.py <nlm_pdf> <section_title> --template weekly.pptx \\
      [--titles titles.txt] [-o output.pptx]

範例：
  python3 nlm_embed_slides.py NLM_B_投影片_20260406.pdf "全球總經與操作策略" \\
      --template Weekly20260406.pptx --titles titles_B.txt -o output_B.pptx

titles.txt 格式（每行一個標題，對應每頁，第一行對應封面會被跳過）：
  封面（跳過）
  核心摘要：五大總經趨勢概覽
  財政赤字隱憂：突破經濟成長率 6% 的尾部風險
  ...

如果不提供 --titles，副標題留空，可事後在 PowerPoint 填入。
"""

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# ── 裁剪參數 ──────────────────────────────────────────────────────
TITLE_BAR_THRESHOLD = 200   # RGB 均值低於此 = 深色標題列
CONTENT_START_THRESHOLD = 230  # RGB 均值高於此 = 內容區開始
WATERMARK_SCAN_COLS = 300   # 從右側掃描浮水印的像素寬度
DPI = 250                   # PDF 轉圖解析度

# ── Layout 常量（底稿 Layout 5「只有標題」）──────────────────────
LAYOUT_NAME = "只有標題"
LAYOUT_INDEX = 5
PH_TITLE_IDX = 0    # 區段標題
PH_SUBTITLE_IDX = 1  # 副標題（NLM 頁面標題）
PH_PICTURE_IDX = 2   # 圖片 placeholder


def pdf_to_images(pdf_path: Path, dpi: int = DPI) -> list[Path]:
    """用 pdftoppm 將 PDF 轉為 PNG"""
    with tempfile.TemporaryDirectory() as tmpdir:
        prefix = Path(tmpdir) / "page"
        subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path), str(prefix)],
            check=True, capture_output=True,
        )
        pages = sorted(Path(tmpdir).glob("page-*.png"))
        out_dir = Path(tempfile.mkdtemp(prefix="nlm_crop_"))
        result = []
        for p in pages:
            dst = out_dir / p.name
            p.rename(dst)
            result.append(dst)
        return result


def detect_crop_bounds(img: Image.Image) -> tuple[int, int]:
    """動態偵測裁剪邊界：標題列結束 y 和浮水印開始 y"""
    pixels = np.array(img)
    w, h = img.size
    sample_left = w // 6
    sample_right = 5 * w // 6

    # 找標題列結束位置
    in_dark = False
    content_start = 0
    for y in range(h // 3):
        row_mean = pixels[y, sample_left:sample_right, :].mean()
        if not in_dark and row_mean < TITLE_BAR_THRESHOLD:
            in_dark = True
        elif in_dark and row_mean > CONTENT_START_THRESHOLD:
            content_start = y
            break

    if content_start == 0:
        content_start = int(h * 0.12)

    # 找浮水印位置（右下角 "※ NotebookLM"）
    watermark_top = h
    for y in range(h - 1, int(h * 0.9), -1):
        right_region = pixels[y, w - WATERMARK_SCAN_COLS:w - 10, :]
        if right_region.mean() < 240:
            watermark_top = y + 5
            break

    if watermark_top >= h:
        watermark_top = int(h * 0.98)

    return content_start, watermark_top


def crop_slide(img: Image.Image, skip_title_detect: bool = False) -> Image.Image:
    """裁剪單頁：去掉標題列和浮水印"""
    w, h = img.size
    if skip_title_detect:
        top, bottom = 0, int(h * 0.98)
    else:
        top, bottom = detect_crop_bounds(img)
    return img.crop((0, top, w, bottom))


def create_pptx(
    template_path: Path,
    cropped_images: list[Path],
    section_title: str,
    subtitles: list[str],
    output_path: Path,
    skip_cover: bool = True,
):
    """用底稿 PPTX 的 layout 建立投影片"""
    prs = Presentation(str(template_path))

    # 找到「只有標題」layout
    layout = None
    for sl in prs.slide_layouts:
        if sl.name == LAYOUT_NAME:
            layout = sl
            break
    if layout is None:
        # fallback: 用 index
        layout = prs.slide_layouts[LAYOUT_INDEX]
        print(f"  [WARN] 找不到 layout '{LAYOUT_NAME}'，用 index {LAYOUT_INDEX}")

    # 移除底稿中的所有現有投影片（只保留 layout）
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    start_idx = 1 if skip_cover else 0

    for i, img_path in enumerate(cropped_images[start_idx:], start=start_idx):
        slide = prs.slides.add_slide(layout)

        # 填入標題 placeholder (idx=0)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == PH_TITLE_IDX:
                ph.text = section_title
            elif ph.placeholder_format.idx == PH_SUBTITLE_IDX:
                subtitle_text = subtitles[i] if i < len(subtitles) else ""
                ph.text = subtitle_text

        # 插入裁剪後的圖片
        # 用 add_picture 而非 picture placeholder（因為尺寸需要動態調整）
        img = Image.open(img_path)
        img_w, img_h = img.size
        aspect = img_h / img_w

        # 圖片全寬，y 從副標題下方開始
        pic_width = Emu(9144000)  # 10"
        pic_height = int(pic_width * aspect)
        pic_top = Emu(2064000)    # ~2.26"

        # 確保不超出投影片底部
        max_height = prs.slide_height - pic_top - Emu(50000)
        if pic_height > max_height:
            pic_height = max_height

        slide.shapes.add_picture(
            str(img_path), Emu(0), pic_top, pic_width, Emu(pic_height)
        )

        # 刪除 picture placeholder（idx=2）避免空白佔位框
        for ph in list(slide.placeholders):
            if ph.placeholder_format.idx == PH_PICTURE_IDX:
                sp = ph._element
                sp.getparent().remove(sp)

    prs.save(str(output_path))
    count = len(cropped_images) - (1 if skip_cover else 0)
    print(f"  [PPTX] 產出 {count} 頁（layout: {layout.name}）")
    print(f"  [PPTX] 儲存至 {output_path}")


def run(
    pdf_path: Path,
    section_title: str,
    template_path: Path,
    titles_file: Path | None = None,
    output_path: Path | None = None,
    skip_cover: bool = True,
):
    """主流程"""
    print(f"{'='*50}")
    print(f"  NLM 投影片裁剪嵌入")
    print(f"{'='*50}")
    print(f"  PDF:    {pdf_path.name}")
    print(f"  底稿:   {template_path.name}")
    print(f"  區段:   {section_title}")

    # 讀取標題列表
    subtitles = []
    if titles_file and titles_file.exists():
        subtitles = [
            line.strip()
            for line in titles_file.read_text(encoding="utf-8").splitlines()
            if line.strip()
        ]
        print(f"  標題:   {len(subtitles)} 個（from {titles_file.name}）")
    else:
        print(f"  標題:   未提供，副標題將留空")

    # PDF → PNG
    print(f"\n  [1/3] PDF → PNG (DPI={DPI}) ...")
    page_images = pdf_to_images(pdf_path)
    print(f"  共 {len(page_images)} 頁")

    # 裁剪
    print(f"  [2/3] 裁剪標題列 + 浮水印 ...")
    cropped_dir = Path(tempfile.mkdtemp(prefix="nlm_embed_"))
    cropped_images = []
    for i, img_path in enumerate(page_images):
        img = Image.open(img_path)
        is_cover = (i == 0)
        cropped = crop_slide(img, skip_title_detect=is_cover)
        out_path = cropped_dir / f"cropped_{i+1:02d}.png"
        cropped.save(out_path, "PNG", optimize=True)
        cropped_images.append(out_path)

        w, h = img.size
        cw, ch = cropped.size
        pct = (1 - ch / h) * 100
        print(f"    P{i+1}: {w}x{h} → {cw}x{ch} (裁掉 {pct:.1f}%)")

    # 建立 PPTX
    print(f"  [3/3] 建立 PPTX ...")
    if output_path is None:
        output_path = pdf_path.parent / f"{pdf_path.stem}_embed.pptx"

    create_pptx(
        template_path, cropped_images, section_title,
        subtitles, output_path, skip_cover,
    )

    # 清理
    import shutil
    for img_path in page_images:
        img_path.unlink(missing_ok=True)
    shutil.rmtree(cropped_dir, ignore_errors=True)

    print(f"\n  完成！")
    return str(output_path)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="裁剪 NLM 投影片並嵌入 PPTX")
    parser.add_argument("pdf", type=Path, help="NLM 投影片 PDF 路徑")
    parser.add_argument("section_title", help="區段標題（如「全球總經與操作策略」）")
    parser.add_argument("--template", type=Path, required=True, help="底稿 PPTX 路徑")
    parser.add_argument("--titles", type=Path, help="標題文字檔（每行一個）")
    parser.add_argument("-o", "--output", type=Path, help="輸出 PPTX 路徑")
    parser.add_argument("--keep-cover", action="store_true", help="保留 NLM 封面頁")
    args = parser.parse_args()

    if not args.pdf.exists():
        print(f"錯誤: PDF 不存在 {args.pdf}")
        sys.exit(1)
    if not args.template.exists():
        print(f"錯誤: 底稿不存在 {args.template}")
        sys.exit(1)

    run(
        pdf_path=args.pdf,
        section_title=args.section_title,
        template_path=args.template,
        titles_file=args.titles,
        output_path=args.output,
        skip_cover=not args.keep_cover,
    )
