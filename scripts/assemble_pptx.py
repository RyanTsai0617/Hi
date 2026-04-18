#!/usr/bin/env python3
"""
NLM Slide Deck → Weekly PPTX Assembler
---------------------------------------
Takes an NLM slide deck PDF, crops each page (remove title + watermark),
and assembles into a Weekly-style PPTX using a template.

Usage:
    python assemble_pptx.py \
        --pdf  nlm_slide_deck.pdf \
        --template Weekly20260406.pptx \
        --output  output.pptx \
        --slides slides_config.json

slides_config.json format:
{
    "cover_title": "利率觀點與策略看法",
    "cover_subtitle": "資料更新日：2026/4/6",
    "pages": [
        {"title": "總經環境與市場概況", "subtitle": "通膨問題 vs. 經濟成長擔憂", "page": 2},
        {"title": "總經環境與市場概況", "subtitle": null, "page": 1},
        ...
    ]
}

If --slides is omitted, all pages are used with generic titles.
"""

import argparse
import json
import os
import sys
import tempfile
from pathlib import Path

# ── Dependencies ───────────────────────────────────────────────────
try:
    import fitz  # PyMuPDF
except ImportError:
    sys.exit("ERROR: pip install PyMuPDF")
try:
    from PIL import Image as PILImage
except ImportError:
    sys.exit("ERROR: pip install Pillow")
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Constants ──────────────────────────────────────────────────────
SLIDE_W = 9144000
SLIDE_H = 6858000
MAX_IMG_W = Cm(25)   # 25 cm
MAX_IMG_H = Cm(14)   # 14 cm
CROP_TOP_PCT = 0.145  # 14.5% from top (NLM title)
CROP_BOT_PCT = 0.035  # 3.5% from bottom (NLM watermark)
FONT = "Microsoft JhengHei"
RENDER_SCALE = 2.5    # PDF rendering scale for crisp images


def pdf_to_cropped_pngs(pdf_path, output_dir):
    """Convert PDF pages to cropped PNGs (remove NLM title + watermark)."""
    os.makedirs(output_dir, exist_ok=True)
    doc = fitz.open(pdf_path)
    paths = []
    for i, page in enumerate(doc):
        mat = fitz.Matrix(RENDER_SCALE, RENDER_SCALE)
        pix = page.get_pixmap(matrix=mat)
        # Save full page first
        full_path = os.path.join(output_dir, f"_full_{i+1:02d}.png")
        pix.save(full_path)
        # Crop
        img = PILImage.open(full_path)
        w, h = img.size
        top = int(h * CROP_TOP_PCT)
        bot = int(h * (1 - CROP_BOT_PCT))
        cropped = img.crop((0, top, w, bot))
        out = os.path.join(output_dir, f"page_{i+1:02d}.png")
        cropped.save(out)
        paths.append(out)
        os.remove(full_path)
        print(f"  Page {i+1}: {w}x{h} → cropped {cropped.size[0]}x{cropped.size[1]}")
    doc.close()
    return paths


def assemble(template_path, png_paths, slides_config, output_path):
    """Assemble cropped PNGs into Weekly-style PPTX."""
    prs = Presentation(template_path)

    # Delete existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    title_layout = prs.slide_layouts[0]   # Cover
    content_layout = prs.slide_layouts[5]  # Title Only

    # ── Cover slide ────────────────────────────────────────────────
    sl = prs.slides.add_slide(title_layout)
    for shape in sl.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = slides_config.get("cover_title", "Weekly Report")
        elif shape.placeholder_format.idx == 1:
            shape.text = slides_config.get("cover_subtitle", "")

    # ── Content slides ─────────────────────────────────────────────
    pages_cfg = slides_config.get("pages", [])

    for i, png_path in enumerate(png_paths):
        sl = prs.slides.add_slide(content_layout)

        # Title from config or default
        if i < len(pages_cfg):
            cfg = pages_cfg[i]
            title = cfg.get("title", f"Slide {i+1}")
            subtitle = cfg.get("subtitle")
        else:
            title = f"Slide {i+1}"
            subtitle = None

        # Set title and subtitle via placeholders
        ph2_removed = False
        for shape in list(sl.placeholders):
            idx = shape.placeholder_format.idx
            if idx == 0:
                shape.text = title
            elif idx == 1:
                if subtitle:
                    shape.text = subtitle
                else:
                    # Remove subtitle placeholder if no subtitle
                    sp = shape._element
                    sp.getparent().remove(sp)
            elif idx == 2:
                # Remove picture placeholder (we'll add image manually for precise control)
                sp = shape._element
                sp.getparent().remove(sp)
                ph2_removed = True

        has_sub = bool(subtitle)

        # Image placement — use template's content area
        # With subtitle: top = 2062079 (below subtitle)
        # Without subtitle: top = 1528004 (below title only)
        if os.path.exists(png_path):
            with PILImage.open(png_path) as img:
                iw, ih = img.size
            aspect = iw / ih

            fit_w = MAX_IMG_W
            fit_h = int(fit_w / aspect)
            if fit_h > MAX_IMG_H:
                fit_h = MAX_IMG_H
                fit_w = int(fit_h * aspect)

            area_top = 2062079 if has_sub else 1528004
            area_bot = SLIDE_H - 200000
            area_h = area_bot - area_top

            left = (SLIDE_W - fit_w) // 2
            top = area_top + (area_h - fit_h) // 2

            sl.shapes.add_picture(png_path, left, top, fit_w, fit_h)

    prs.save(output_path)
    print(f"\nPPTX saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


def main():
    parser = argparse.ArgumentParser(description="NLM Slide Deck → Weekly PPTX")
    parser.add_argument("--pdf", required=True, help="NLM slide deck PDF path")
    parser.add_argument("--template", required=True, help="Weekly PPTX template path")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    parser.add_argument("--slides", default=None, help="Slides config JSON (optional)")
    parser.add_argument("--crop-top", type=float, default=CROP_TOP_PCT,
                        help="Top crop percentage (default: 0.145)")
    parser.add_argument("--crop-bot", type=float, default=CROP_BOT_PCT,
                        help="Bottom crop percentage (default: 0.035)")
    parser.add_argument("--max-w", type=float, default=25,
                        help="Max image width in cm (default: 25)")
    parser.add_argument("--max-h", type=float, default=14,
                        help="Max image height in cm (default: 14)")
    args = parser.parse_args()

    # Update module-level constants from CLI args
    import sys
    mod = sys.modules[__name__]
    mod.CROP_TOP_PCT = args.crop_top
    mod.CROP_BOT_PCT = args.crop_bot
    mod.MAX_IMG_W = Cm(args.max_w)
    mod.MAX_IMG_H = Cm(args.max_h)

    if not Path(args.pdf).exists():
        sys.exit(f"ERROR: PDF not found: {args.pdf}")
    if not Path(args.template).exists():
        sys.exit(f"ERROR: Template not found: {args.template}")

    # Slides config
    if args.slides and Path(args.slides).exists():
        with open(args.slides, "r", encoding="utf-8") as f:
            slides_config = json.load(f)
    else:
        slides_config = {
            "cover_title": "Weekly Report",
            "cover_subtitle": "",
            "pages": []
        }

    # Step 1: PDF → cropped PNGs
    tmp_dir = os.path.join(os.path.dirname(args.output), "_nlm_cropped")
    print(f"Converting PDF → cropped PNGs (crop top={CROP_TOP_PCT:.1%}, bot={CROP_BOT_PCT:.1%})...")
    png_paths = pdf_to_cropped_pngs(args.pdf, tmp_dir)
    print(f"  {len(png_paths)} pages exported\n")

    # Step 2: Assemble PPTX
    print(f"Assembling PPTX (max image: {args.max_w}cm x {args.max_h}cm)...")
    assemble(args.template, png_paths, slides_config, args.output)


if __name__ == "__main__":
    main()
