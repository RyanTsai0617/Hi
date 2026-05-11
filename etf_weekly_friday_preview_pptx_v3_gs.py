"""ETF Weekly Friday topic preview — v3 Goldman Sachs research style.

Visual: deep navy banner + mustard gold accents + dense tables + "BOTTOM LINE" callouts.
9 slides, 16:9. Sharp rectangles, no rounded corners. High info density.

vs v2:
- 配色：teal → GS navy (#00295C) + gold (#FFB81C) + black/white/gray
- 字體：粗體標題 + 細體內文、英中對照銜接
- Banner：頂部 GS_BLUE 0.35" 高、左 "EQUITY RESEARCH | ETF WEEKLY" 白字、右頁碼
- Cover：黑大標 + 3 KPI 矩形卡（GS_BLUE / GS_GOLD / 黑色）、無圓角
- Section divider：黑底 + 巨大 GS_GOLD 編號 + 白色標題
- Topic cards：白底細邊框、左 GS_GOLD「01」編號、底「BOTTOM LINE」黃線
- ETF matrix：GS_BLUE header + 灰白交錯 + heatmap 用 navy/red 漸層
- Macro：左右雙欄 + 細金線分隔 + 編號清單
- Monday tasks：8 行緊湊清單、左 GS_GOLD 編號圓圈

本地測試用 0511 期 mock data、產 PPTX。
"""
from __future__ import annotations
import os
from pathlib import Path
from io import BytesIO

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ============== Goldman Sachs visual palette ==============
GS_BLUE = RGBColor(0x00, 0x29, 0x5C)         # Goldman navy
GS_BLUE_LIGHT = RGBColor(0x1F, 0x4E, 0x82)   # lighter navy (heatmap)
GS_GOLD = RGBColor(0xFF, 0xB8, 0x1C)         # mustard
GS_GOLD_DARK = RGBColor(0xC8, 0x8E, 0x10)    # darker gold for text on white
BLACK = RGBColor(0x10, 0x10, 0x10)
DARK_GRAY = RGBColor(0x2A, 0x2A, 0x2A)
MED_GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
ALT_ROW = RGBColor(0xF5, 0xF5, 0xF5)
BORDER_GRAY = RGBColor(0xC8, 0xC8, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Loss/gain — conservative GS palette
GAIN = RGBColor(0x1A, 0x6E, 0x1A)
LOSS = RGBColor(0xA5, 0x0B, 0x0B)

# Card backgrounds (for KPI strip on cover)
CARD_BG_1 = GS_BLUE         # primary
CARD_BG_2 = GS_GOLD         # accent
CARD_BG_3 = DARK_GRAY       # neutral

# Hex versions for matplotlib
GS_BLUE_HEX = "#00295C"
GS_BLUE_LIGHT_HEX = "#1F4E82"
GS_GOLD_HEX = "#FFB81C"
GAIN_HEX = "#1A6E1A"
LOSS_HEX = "#A50B0B"
DARK_GRAY_HEX = "#2A2A2A"
MED_GRAY_HEX = "#555555"
LIGHT_GRAY_HEX = "#E8E8E8"
WHITE_HEX = "#FFFFFF"

ZH_FONT = "微軟正黑體"
EN_FONT = "Calibri"

for fpath in ["C:/Windows/Fonts/msjh.ttc", "C:/Windows/Fonts/msjh.ttf",
              "/Library/Fonts/Microsoft Sans Serif.ttf"]:
    if Path(fpath).exists():
        font_manager.fontManager.addfont(fpath)
        break
plt.rcParams["font.family"] = ["Microsoft JhengHei", "Microsoft YaHei", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False


# ============== python-pptx helpers ==============

def apply_font(run, *, size=14, bold=False, color=DARK_GRAY, font_en=EN_FONT, font_zh=ZH_FONT, italic=False):
    run.font.name = font_en
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_zh)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    paras = text.split("\n") if isinstance(text, str) else text
    for i, p_text in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = p_text
        apply_font(run, size=size, bold=bold, color=color, italic=italic)
    return tb


def add_rect(slide, x, y, w, h, fill, *, line_color=None, line_weight=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_weight)
    return s


def add_line(slide, x1, y1, x2, y2, color, weight=0.75):
    """Thin horizontal/vertical line."""
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_image(slide, x, y, w, h, image_stream):
    return slide.shapes.add_picture(image_stream, x, y, width=w, height=h)


# ============== Page chrome (banner + footer) ==============

GS_HEADER_TEXT = "EQUITY RESEARCH  |  ETF WEEKLY"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_page_chrome(slide, *, page_num, total, footer_left, footer_right=None):
    """Top GS_BLUE banner + bottom gold line + footer."""
    # Top banner
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.32), GS_BLUE)
    # Banner text left
    add_text(slide, Inches(0.4), Inches(0.04), Inches(8.0), Inches(0.24),
             GS_HEADER_TEXT, size=10, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # Page number top right (small white)
    add_text(slide, Inches(11.5), Inches(0.04), Inches(1.6), Inches(0.24),
             f"PAGE {page_num} / {total}",
             size=9, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # Gold line under banner
    add_rect(slide, Inches(0), Inches(0.32), SLIDE_W, Inches(0.04), GS_GOLD)

    # Bottom gold line + footer
    add_rect(slide, Inches(0.4), Inches(7.18), Inches(12.5), Inches(0.02), GS_GOLD)
    add_text(slide, Inches(0.4), Inches(7.22), Inches(8.0), Inches(0.22),
             footer_left, size=8, color=MED_GRAY, anchor=MSO_ANCHOR.MIDDLE)
    if footer_right:
        add_text(slide, Inches(8.4), Inches(7.22), Inches(4.5), Inches(0.22),
                 footer_right, size=8, color=MED_GRAY, italic=True,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_section_title(slide, *, title_en, title_zh, subtitle=None, y=Inches(0.55)):
    """Page-level title block: small EN cap + big ZH + optional subtitle."""
    # EN small uppercase
    add_text(slide, Inches(0.5), y, Inches(12.5), Inches(0.25),
             title_en, size=10, bold=True, color=GS_GOLD_DARK,
             anchor=MSO_ANCHOR.TOP)
    # ZH big bold
    add_text(slide, Inches(0.5), y + Inches(0.28), Inches(12.5), Inches(0.55),
             title_zh, size=22, bold=True, color=BLACK,
             anchor=MSO_ANCHOR.TOP)
    if subtitle:
        add_text(slide, Inches(0.5), y + Inches(0.85), Inches(12.5), Inches(0.4),
                 subtitle, size=11, color=MED_GRAY, italic=True,
                 anchor=MSO_ANCHOR.TOP)
    # thin gray line under title block
    line_y = y + (Inches(1.30) if subtitle else Inches(0.92))
    add_rect(slide, Inches(0.5), line_y, Inches(12.3), Inches(0.01), BORDER_GRAY)


# ============== Matplotlib charts (GS palette) ==============

def make_reversal_bar_chart(rebound_rows, decline_rows, *, width_in=12.0, height_in=4.3):
    """Two-panel horizontal bar — rebound left (GS_BLUE) + decline right (LOSS)."""
    fig, axes = plt.subplots(1, 2, figsize=(width_in, height_in), dpi=160)
    fig.patch.set_facecolor(WHITE_HEX)

    # Left: rebound
    ax = axes[0]
    rebound = sorted(rebound_rows, key=lambda r: r[3])
    labels = [r[0] for r in rebound]
    values = [r[3] for r in rebound]
    bars = ax.barh(labels, values, color=GS_BLUE_HEX, edgecolor="none", height=0.72)
    for bar, v in zip(bars, values):
        ax.text(v + max(values) * 0.02, bar.get_y() + bar.get_height() / 2,
                f"+{v:.1f}", va="center", fontsize=9, color=GS_BLUE_HEX, weight="bold")
    ax.set_title("REBOUND  /  本週反彈", fontsize=11, color=DARK_GRAY_HEX,
                 pad=10, weight="bold", loc="left")
    ax.set_xlabel("reversal magnitude (pp)", fontsize=8, color=MED_GRAY_HEX)
    ax.set_facecolor(WHITE_HEX)
    for spine in ax.spines.values(): spine.set_visible(False)
    ax.spines["bottom"].set_visible(True)
    ax.spines["bottom"].set_color(BORDER_GRAY.__class__.__name__ if False else "#C8C8C8")
    ax.tick_params(colors=MED_GRAY_HEX, labelsize=9)
    ax.grid(axis="x", linestyle=":", alpha=0.25, color=MED_GRAY_HEX)
    ax.set_axisbelow(True)

    # Right: decline
    ax = axes[1]
    decline = sorted(decline_rows, key=lambda r: r[3], reverse=True)
    labels = [r[0] for r in decline]
    values = [r[3] for r in decline]
    bars = ax.barh(labels, values, color=LOSS_HEX, edgecolor="none", height=0.72)
    for bar, v in zip(bars, values):
        ax.text(v - abs(min(values)) * 0.02, bar.get_y() + bar.get_height() / 2,
                f"{v:.1f}", va="center", ha="right", fontsize=9, color=LOSS_HEX, weight="bold")
    ax.set_title("DECLINE  /  本週反向下殺", fontsize=11, color=DARK_GRAY_HEX,
                 pad=10, weight="bold", loc="left")
    ax.set_xlabel("reversal magnitude (pp)", fontsize=8, color=MED_GRAY_HEX)
    ax.set_facecolor(WHITE_HEX)
    for spine in ax.spines.values(): spine.set_visible(False)
    ax.spines["bottom"].set_visible(True)
    ax.tick_params(colors=MED_GRAY_HEX, labelsize=9)
    ax.grid(axis="x", linestyle=":", alpha=0.25, color=MED_GRAY_HEX)
    ax.set_axisbelow(True)

    plt.tight_layout()
    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=160, bbox_inches="tight",
                facecolor=WHITE_HEX, edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


# ============== Slide builders ==============

def build_slide_cover(slide, *, period, week_range, status_text, kpi_cards, footer_left, footer_right):
    """Cover slide: GS-style with banner, big black title, 3 sharp KPI rectangles."""
    # Top GS_BLUE banner with main brand line
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.55), GS_BLUE)
    add_text(slide, Inches(0.5), Inches(0.10), Inches(9.0), Inches(0.36),
             "EQUITY RESEARCH  |  ETF WEEKLY  |  FRIDAY TOPIC PREVIEW",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(9.5), Inches(0.10), Inches(3.5), Inches(0.36),
             "CATHAY SECURITIES",
             size=11, bold=True, color=GS_GOLD,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0), Inches(0.55), SLIDE_W, Inches(0.04), GS_GOLD)

    # Big title
    add_text(slide, Inches(0.5), Inches(1.1), Inches(12.5), Inches(0.55),
             "U.S. ETF WEEKLY", size=14, bold=True, color=GS_GOLD_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.5), Inches(1.55), Inches(12.5), Inches(1.0),
             "美股 ETF 週報",
             size=46, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.5), Inches(2.55), Inches(12.5), Inches(0.5),
             "Friday Topic Preview · 議題候選掃描",
             size=18, color=DARK_GRAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Period meta (centered)
    add_text(slide, Inches(0.5), Inches(3.1), Inches(12.5), Inches(0.4),
             f"PERIOD  {period}     |     WEEK  {week_range}",
             size=11, bold=True, color=GS_BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.5), Inches(3.5), Inches(12.5), Inches(0.4),
             status_text,
             size=10, color=MED_GRAY, italic=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 3 sharp KPI cards (no rounded corners)
    card_w = Inches(3.7)
    card_h = Inches(1.8)
    gap = Inches(0.25)
    cards_y = Inches(4.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2

    for i, (ticker, value, label, bg_color) in enumerate(kpi_cards):
        cx = start_x + i * (card_w + gap)
        # Sharp rectangle (no rounded)
        add_rect(slide, cx, cards_y, card_w, card_h, bg_color)
        # Top: small label "TOP REBOUND" / "TOP DECLINE" / "WATCH" type tag
        tag_text = ["TOP REBOUND", "TOP DECLINE", "NOTABLE"][i] if i < 3 else "ITEM"
        add_text(slide, cx + Inches(0.15), cards_y + Inches(0.1),
                 card_w - Inches(0.3), Inches(0.25), tag_text,
                 size=9, bold=True,
                 color=(GS_GOLD if bg_color == GS_BLUE else (BLACK if bg_color == GS_GOLD else GS_GOLD)),
                 anchor=MSO_ANCHOR.MIDDLE)
        # Ticker
        text_color = WHITE if bg_color != GS_GOLD else BLACK
        add_text(slide, cx, cards_y + Inches(0.4), card_w, Inches(0.4), ticker,
                 size=18, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Value (big)
        add_text(slide, cx, cards_y + Inches(0.78), card_w, Inches(0.7), value,
                 size=32, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Label
        add_text(slide, cx + Inches(0.15), cards_y + Inches(1.45),
                 card_w - Inches(0.3), Inches(0.3), label,
                 size=10, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom analyst block
    add_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.02), GS_GOLD)
    add_text(slide, Inches(0.5), Inches(6.55), Inches(8.0), Inches(0.3),
             "Ryan Tsai  |  法人債券業務部 · 法人業務二處",
             size=10, bold=True, color=DARK_GRAY)
    add_text(slide, Inches(0.5), Inches(6.85), Inches(8.0), Inches(0.3),
             "tsaitsungyi0617@gmail.com  |  Cathay Securities Institutional Bond Dept.",
             size=9, color=MED_GRAY)
    add_text(slide, Inches(8.5), Inches(6.55), Inches(4.3), Inches(0.6),
             "Important Disclosures: Internal use only.\nNot for distribution.",
             size=8, color=MED_GRAY, italic=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)

    # Footer
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8.0), Inches(0.22),
             footer_left, size=8, color=MED_GRAY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(8.5), Inches(7.22), Inches(4.3), Inches(0.22),
             footer_right or "PAGE 1 / 9", size=8, color=MED_GRAY,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def build_slide_section_divider(slide, *, section_num, section_title_en, section_title_zh, footer_left):
    """Black bg + huge gold section number + white title."""
    # Full black bg
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BLACK)
    # Top brand bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), GS_GOLD)

    # Big gold section number left
    add_text(slide, Inches(0.6), Inches(2.0), Inches(5.0), Inches(3.0),
             section_num, size=200, bold=True, color=GS_GOLD,
             anchor=MSO_ANCHOR.MIDDLE)

    # Vertical thin gold line
    add_rect(slide, Inches(5.7), Inches(2.6), Inches(0.02), Inches(2.0), GS_GOLD)

    # Section title EN (small) + ZH (big)
    add_text(slide, Inches(6.0), Inches(2.7), Inches(7.0), Inches(0.4),
             section_title_en, size=14, bold=True, color=GS_GOLD,
             anchor=MSO_ANCHOR.TOP)
    add_text(slide, Inches(6.0), Inches(3.2), Inches(7.0), Inches(1.5),
             section_title_zh, size=42, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.TOP)

    # Bottom footer
    add_rect(slide, Inches(0.5), Inches(7.18), Inches(12.3), Inches(0.02), GS_GOLD)
    add_text(slide, Inches(0.5), Inches(7.22), Inches(12.3), Inches(0.22),
             footer_left, size=8, color=LIGHT_GRAY, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)


def build_slide_reversal_chart(slide, *, rebound_rows, decline_rows, page_num, total, footer_left):
    add_page_chrome(slide, page_num=page_num, total=total, footer_left=footer_left)
    add_section_title(slide,
                      title_en="REVERSAL SIGNALS",
                      title_zh="反轉訊號 — 本週 vs 上週",
                      subtitle="Top 10 by reversal magnitude · GS Blue = rebound, Red = decline")

    # BOTTOM LINE callout (light gray bg + gold left bar)
    callout_y = Inches(2.05)
    callout_h = Inches(0.55)
    add_rect(slide, Inches(0.5), callout_y, Inches(0.08), callout_h, GS_GOLD)
    add_rect(slide, Inches(0.58), callout_y, Inches(12.22), callout_h, ALT_ROW)
    top_rebound = max(rebound_rows, key=lambda r: r[3])
    top_decline = min(decline_rows, key=lambda r: r[3])
    callout = (f"BOTTOM LINE   {top_rebound[0]} {top_rebound[3]:+.1f}pp 最強反彈"
               f"（本週 {top_rebound[1]:+.2f}% / 上週 {top_rebound[2]:+.2f}%）  ·  "
               f"{top_decline[0]} {top_decline[3]:+.1f}pp 最猛下殺"
               f"（本週 {top_decline[1]:+.2f}% / 上週 {top_decline[2]:+.2f}%）")
    add_text(slide, Inches(0.75), callout_y, Inches(11.9), callout_h,
             callout, size=11, bold=True, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

    # Chart
    chart_buf = make_reversal_bar_chart(rebound_rows, decline_rows,
                                         width_in=12.0, height_in=4.3)
    add_image(slide, Inches(0.5), Inches(2.75), Inches(12.3), Inches(4.2), chart_buf)


def build_slide_topic_cards(slide, *, title_en, title_zh, subtitle,
                             candidates, page_num, total, footer_left,
                             accent_label="THEME"):
    """Investment themes — sharp white cards with gold left number, thin border, BOTTOM LINE row."""
    add_page_chrome(slide, page_num=page_num, total=total, footer_left=footer_left)
    add_section_title(slide,
                      title_en=title_en,
                      title_zh=title_zh,
                      subtitle=subtitle)

    n = len(candidates)
    available_h = Inches(5.0)
    gap = Inches(0.15)
    card_h = (available_h - gap * (n - 1)) / n
    card_y_start = Inches(2.05)
    card_w = Inches(12.3)
    card_x = Inches(0.5)

    for i, c in enumerate(candidates):
        cy = card_y_start + i * (card_h + gap)
        # Card bg white + thin gray border
        add_rect(slide, card_x, cy, card_w, card_h, WHITE,
                 line_color=BORDER_GRAY, line_weight=0.75)
        # Left dark blue number box
        num_w = Inches(0.85)
        add_rect(slide, card_x, cy, num_w, card_h, GS_BLUE)
        add_text(slide, card_x, cy + Inches(0.05), num_w, Inches(0.22),
                 accent_label, size=7.5, bold=True, color=GS_GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Big number — adaptive size
        num_size = 30 if card_h >= Inches(1.3) else 22
        add_text(slide, card_x, cy + Inches(0.25), num_w, card_h - Inches(0.3),
                 f"{i+1:02d}", size=num_size, bold=True, color=GS_GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Card content (right of number)
        content_x = card_x + num_w + Inches(0.2)
        content_w = card_w - num_w - Inches(0.4)

        is_tall = card_h >= Inches(1.3)

        if is_tall:
            # Full 4-zone layout for main cards (3 cards × ~1.57" each)
            # Zone 1: Headline (top)
            add_text(slide, content_x, cy + Inches(0.08), content_w, Inches(0.34),
                     c.get("主標", ""), size=13, bold=True, color=BLACK,
                     anchor=MSO_ANCHOR.TOP)
            # Zone 2: KEY DATA + ETFs (two-column inline)
            stats_y = cy + Inches(0.5)
            add_text(slide, content_x, stats_y, Inches(1.0), Inches(0.24),
                     "KEY DATA", size=8, bold=True, color=GS_GOLD_DARK,
                     anchor=MSO_ANCHOR.TOP)
            add_text(slide, content_x + Inches(1.0), stats_y, Inches(5.5), Inches(0.24),
                     c.get("核心數據", ""), size=10, bold=True, color=DARK_GRAY,
                     anchor=MSO_ANCHOR.TOP)
            add_text(slide, content_x + Inches(6.5), stats_y, Inches(0.6), Inches(0.24),
                     "ETFs", size=8, bold=True, color=GS_GOLD_DARK,
                     anchor=MSO_ANCHOR.TOP)
            add_text(slide, content_x + Inches(7.1), stats_y, content_w - Inches(7.1), Inches(0.24),
                     c.get("相關 ETF", ""), size=9, color=DARK_GRAY,
                     anchor=MSO_ANCHOR.TOP)
            # Zone 3: Driver
            drv_y = cy + Inches(0.78)
            add_text(slide, content_x, drv_y, Inches(1.0), Inches(0.24),
                     "DRIVER", size=8, bold=True, color=GS_GOLD_DARK,
                     anchor=MSO_ANCHOR.TOP)
            add_text(slide, content_x + Inches(1.0), drv_y, content_w - Inches(1.0), Inches(0.24),
                     c.get("driver", ""), size=9.5, color=DARK_GRAY,
                     anchor=MSO_ANCHOR.TOP)
            # Zone 4: Thin gold separator + BOTTOM LINE at very bottom
            sep_y = cy + card_h - Inches(0.42)
            add_rect(slide, content_x, sep_y, content_w, Inches(0.01), GS_GOLD)
            bl_y = sep_y + Inches(0.06)
            add_text(slide, content_x, bl_y, Inches(1.5), Inches(0.3),
                     "BOTTOM LINE", size=8, bold=True, color=GS_BLUE,
                     anchor=MSO_ANCHOR.TOP)
            add_text(slide, content_x + Inches(1.5), bl_y, content_w - Inches(1.5), Inches(0.3),
                     c.get("對台灣機構意義", ""), size=10, bold=True, color=BLACK,
                     anchor=MSO_ANCHOR.TOP)
        else:
            # Compressed 3-zone layout for sub cards (5 cards × ~0.88" each)
            # Zone 1: Headline
            add_text(slide, content_x, cy + Inches(0.04), content_w, Inches(0.28),
                     c.get("主標", ""), size=11.5, bold=True, color=BLACK,
                     anchor=MSO_ANCHOR.TOP)
            # Zone 2: KEY DATA inline + ETFs (compact)
            stats = f"{c.get('核心數據', '')}    ·    ETF：{c.get('相關 ETF', '')}"
            add_text(slide, content_x, cy + Inches(0.32), content_w, Inches(0.22),
                     stats, size=8.5, color=DARK_GRAY,
                     anchor=MSO_ANCHOR.TOP)
            # Zone 3: BOTTOM LINE one-liner with thin gold sep
            sep_y = cy + Inches(0.56)
            add_rect(slide, content_x, sep_y, content_w, Inches(0.008), GS_GOLD)
            bl_y = sep_y + Inches(0.04)
            add_text(slide, content_x, bl_y, Inches(1.4), Inches(0.26),
                     "BOTTOM LINE", size=7.5, bold=True, color=GS_BLUE,
                     anchor=MSO_ANCHOR.TOP)
            add_text(slide, content_x + Inches(1.4), bl_y, content_w - Inches(1.4), Inches(0.26),
                     c.get("對台灣機構意義", ""), size=9, bold=True, color=BLACK,
                     anchor=MSO_ANCHOR.TOP)


def build_slide_etf_matrix(slide, *, theme_name, theme_rows, page_num, total, footer_left):
    add_page_chrome(slide, page_num=page_num, total=total, footer_left=footer_left)
    add_section_title(slide,
                      title_en=f"ETF UNIVERSE  |  {theme_name.upper()}",
                      title_zh=f"ETF 矩陣 — {theme_name}",
                      subtitle="一週% / 上週% / YTD% / 角色定位 · heatmap on YTD column")

    headers = ["TICKER", "NAME / 全名", "WEEK %", "PRIOR WEEK %", "YTD %", "ROLE / 角色定位"]
    n_cols = len(headers)
    n_rows = 1 + len(theme_rows)
    tbl_x = Inches(0.5)
    tbl_y = Inches(2.05)
    tbl_w = Inches(12.3)
    tbl_h = Inches(5.0)

    table = slide.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y, tbl_w, tbl_h).table

    # Set column widths (proportional)
    col_widths = [0.85, 3.2, 1.55, 1.6, 1.55, 3.55]  # sums to 12.3
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)

    # Header
    for j, hd in enumerate(headers):
        c = table.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = GS_BLUE
        c.text_frame.clear()
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hd
        apply_font(r, size=10, bold=True, color=WHITE)

    # Body
    import re
    for i, row in enumerate(theme_rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            t = str(val).strip()
            # Heatmap on WEEK% (col 2) and YTD% (col 4)
            if j in (2, 4) and re.match(r"^[\-+]?\d", t):
                pct_val = float(t.replace("%", "").replace("+", ""))
                if pct_val > 0:
                    intensity = min(abs(pct_val) / 60.0, 1.0) if j == 4 else min(abs(pct_val) / 8.0, 1.0)
                    # Light navy tint (white → GS_BLUE_LIGHT_HEX)
                    bg = RGBColor(
                        max(0xFF - int(0xC0 * intensity), 0x1F),
                        max(0xFF - int(0xB1 * intensity), 0x4E),
                        max(0xFF - int(0x7D * intensity), 0x82),
                    )
                    c.fill.solid()
                    c.fill.fore_color.rgb = bg
                elif pct_val < 0:
                    intensity = min(abs(pct_val) / 60.0, 1.0) if j == 4 else min(abs(pct_val) / 8.0, 1.0)
                    bg = RGBColor(
                        0xFF,
                        max(0xFF - int(0xF4 * intensity), 0x0B),
                        max(0xFF - int(0xF4 * intensity), 0x0B),
                    )
                    c.fill.solid()
                    c.fill.fore_color.rgb = bg
                else:
                    if i % 2 == 0:
                        c.fill.solid()
                        c.fill.fore_color.rgb = ALT_ROW
            elif i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = ALT_ROW

            c.text_frame.clear()
            p = c.text_frame.paragraphs[0]
            if j == 0:  # ticker
                p.alignment = PP_ALIGN.LEFT
            elif j in (2, 3, 4) or re.match(r"^[\-+\$]?\d", t) or t.endswith("%"):
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = t
            color = DARK_GRAY
            bold = (j == 0)
            if t.startswith("+") and re.match(r"^\+\d", t):
                color = GAIN
                bold = True
            elif t.startswith("-") and re.match(r"^-\d", t):
                color = LOSS
                bold = True
            apply_font(r, size=9.5, color=color, bold=bold)


def build_slide_macro_research(slide, *, macro_events, research_reports,
                                page_num, total, footer_left):
    add_page_chrome(slide, page_num=page_num, total=total, footer_left=footer_left)
    add_section_title(slide,
                      title_en="MARKET CATALYSTS & RESEARCH",
                      title_zh="總體事件 + A 級 research",
                      subtitle="本週主要 catalyst · A-grade research notes")

    col_left_x = Inches(0.5)
    col_right_x = Inches(7.0)
    col_w = Inches(6.0)
    col_y = Inches(2.05)
    col_h = Inches(5.0)

    # Left column header
    add_rect(slide, col_left_x, col_y, col_w, Inches(0.4), GS_BLUE)
    add_text(slide, col_left_x + Inches(0.2), col_y, col_w - Inches(0.4), Inches(0.4),
             "MARKET CATALYSTS  /  總體事件",
             size=10, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # Left list — numbered items with thin gold separator
    item_y = col_y + Inches(0.55)
    for i, e in enumerate(macro_events):
        # Number tag
        add_text(slide, col_left_x, item_y, Inches(0.45), Inches(0.45),
                 f"{i+1:02d}", size=12, bold=True, color=GS_GOLD_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # Event text
        add_text(slide, col_left_x + Inches(0.5), item_y,
                 col_w - Inches(0.5), Inches(0.7),
                 e, size=10, color=DARK_GRAY, anchor=MSO_ANCHOR.TOP)
        # Thin gold separator
        item_y_next = item_y + Inches(0.55)
        if i < len(macro_events) - 1:
            add_rect(slide, col_left_x + Inches(0.05), item_y_next - Inches(0.05),
                     col_w - Inches(0.1), Inches(0.005), BORDER_GRAY)
        item_y = item_y_next

    # Right column header
    add_rect(slide, col_right_x, col_y, col_w, Inches(0.4), GS_BLUE)
    add_text(slide, col_right_x + Inches(0.2), col_y, col_w - Inches(0.4), Inches(0.4),
             "A-GRADE RESEARCH  /  研究報告",
             size=10, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # Right list
    if research_reports:
        item_y = col_y + Inches(0.55)
        for i, r in enumerate(research_reports):
            add_text(slide, col_right_x, item_y, Inches(0.45), Inches(0.45),
                     f"{i+1:02d}", size=12, bold=True, color=GS_GOLD_DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
            add_text(slide, col_right_x + Inches(0.5), item_y,
                     col_w - Inches(0.5), Inches(0.7),
                     r, size=10, color=DARK_GRAY, anchor=MSO_ANCHOR.TOP)
            item_y_next = item_y + Inches(0.55)
            if i < len(research_reports) - 1:
                add_rect(slide, col_right_x + Inches(0.05), item_y_next - Inches(0.05),
                         col_w - Inches(0.1), Inches(0.005), BORDER_GRAY)
            item_y = item_y_next
    else:
        add_text(slide, col_right_x + Inches(0.2), col_y + Inches(0.6),
                 col_w - Inches(0.4), Inches(1.0),
                 "(No A-grade research highlights this week.\n本週無重要 A 級 research highlights、待 Monday 補充。)",
                 size=10, color=MED_GRAY, italic=True, anchor=MSO_ANCHOR.TOP)


def build_slide_monday_tasks(slide, *, period, tasks, page_num, total, footer_left):
    add_page_chrome(slide, page_num=page_num, total=total, footer_left=footer_left)
    add_section_title(slide,
                      title_en=f"MONDAY {period} EXECUTION CHECKLIST",
                      title_zh="週一執行清單",
                      subtitle="Estimated 2-2.5 hr full Phase 1-6 · 議題確認 → 重抓 → NLM saturation → docx → Path A/B/C")

    n = len(tasks)
    item_h = Inches(0.52)
    start_y = Inches(2.15)
    row_w = Inches(12.3)
    row_x = Inches(0.5)

    for i, task in enumerate(tasks):
        ty = start_y + i * (item_h + Inches(0.06))
        # Light row bg (alt)
        if i % 2 == 0:
            add_rect(slide, row_x, ty, row_w, item_h, ALT_ROW)
        # Number circle (gold filled)
        num_size = Inches(0.42)
        circle_x = row_x + Inches(0.1)
        circle_y = ty + (item_h - num_size) / 2
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_x, circle_y, num_size, num_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = GS_BLUE
        circle.line.color.rgb = GS_GOLD
        circle.line.width = Pt(1.0)
        # Number in circle
        add_text(slide, circle_x, circle_y, num_size, num_size,
                 f"{i+1}", size=12, bold=True, color=GS_GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Task text
        add_text(slide, row_x + Inches(0.7), ty, row_w - Inches(0.85), item_h,
                 task, size=11, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)


# ============== Main builder ==============

def build_friday_preview(*,
                          this_friday: str,
                          period: str,
                          week_range: str,
                          status_text: str,
                          kpi_cards: list,
                          rebound_rows: list,
                          decline_rows: list,
                          main_candidates: list,
                          sub_candidates: list,
                          theme_tables: list,
                          macro_events: list,
                          research_reports: list,
                          monday_tasks: list,
                          out_path: str):
    """Build the v3 Goldman Sachs style Friday preview PPTX."""
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H

    blank = pres.slide_layouts[6]
    footer_left = f"{this_friday}  |  Cathay Securities Institutional Bond Dept.  |  法人業務二處"

    total = 9

    # 1. Cover
    s1 = pres.slides.add_slide(blank)
    build_slide_cover(s1, period=period, week_range=week_range,
                      status_text=status_text, kpi_cards=kpi_cards,
                      footer_left=footer_left, footer_right="PAGE 1 / 9")

    # 2. Section divider — Reversal
    s2 = pres.slides.add_slide(blank)
    build_slide_section_divider(s2, section_num="01",
                                 section_title_en="REVERSAL SIGNALS",
                                 section_title_zh="反轉訊號",
                                 footer_left=footer_left)

    # 3. Reversal chart
    s3 = pres.slides.add_slide(blank)
    build_slide_reversal_chart(s3, rebound_rows=rebound_rows,
                                decline_rows=decline_rows, page_num=3,
                                total=total, footer_left=footer_left)

    # 4. Investment themes (main 3)
    s4 = pres.slides.add_slide(blank)
    build_slide_topic_cards(s4,
                             title_en="INVESTMENT THEMES — TOP 3 PRIMARY",
                             title_zh="主議題候選 · Top 3",
                             subtitle="敘事最強、優先選為主議題 · Highest conviction narratives",
                             candidates=main_candidates, page_num=4,
                             total=total, footer_left=footer_left,
                             accent_label="MAIN")

    # 5. Investment themes (sub 5)
    s5 = pres.slides.add_slide(blank)
    build_slide_topic_cards(s5,
                             title_en="INVESTMENT THEMES — TOP 5 SECONDARY",
                             title_zh="補充議題候選 · Top 5",
                             subtitle="補充視角、選 2 個搭配主議題 · Pair with primary theme",
                             candidates=sub_candidates, page_num=5,
                             total=total, footer_left=footer_left,
                             accent_label="SUB")

    # 6-7. ETF matrix
    page_num = 6
    for theme_name, theme_rows in theme_tables[:2]:
        s = pres.slides.add_slide(blank)
        build_slide_etf_matrix(s, theme_name=theme_name, theme_rows=theme_rows,
                               page_num=page_num, total=total,
                               footer_left=footer_left)
        page_num += 1

    # 8. Macro + research
    s8 = pres.slides.add_slide(blank)
    build_slide_macro_research(s8, macro_events=macro_events,
                                research_reports=research_reports,
                                page_num=8, total=total,
                                footer_left=footer_left)

    # 9. Monday tasks
    s9 = pres.slides.add_slide(blank)
    build_slide_monday_tasks(s9, period=period, tasks=monday_tasks,
                              page_num=9, total=total, footer_left=footer_left)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    pres.save(out_path)
    print(f"[OK] v3 GS-style Friday preview saved: {out_path}")
    return out_path


# ============== Mock data for local test (0511 期) ==============

if __name__ == "__main__":
    out_path = r"C:/Users/tsait/OneDrive/桌面/每周報告/20260511/_v3_friday_preview_gs.pptx"

    build_friday_preview(
        this_friday="2026-05-08",
        period="20260511",
        week_range="2026-05-04 ~ 2026-05-08",
        status_text="v0 partial scan · Mon-Thu close (4 trading days) · Monday rerun w/ Friday close",
        kpi_cards=[
            ("IRBO", "+7.99%", "本週最強 · AI Brain 鍊條 reignite", CARD_BG_1),
            ("USO", "-5.48%", "YTD 跌穿百分百 · 油市急殺", CARD_BG_2),
            ("GDX", "+5.29%", "V 字反彈 · 黃金回神", CARD_BG_3),
        ],
        rebound_rows=[
            ("GDX", 5.29, -7.66, 12.95),
            ("SLV", 4.85, -0.73, 5.58),
            ("SMH", 5.94, 0.67, 5.27),
            ("SOXX", 5.71, 0.90, 4.81),
            ("XBI", 2.48, -2.15, 4.63),
            ("GLD", 2.01, -2.32, 4.33),
            ("XLK", 4.83, 1.03, 3.80),
            ("EEM", 3.84, 0.61, 3.23),
            ("IBB", 1.69, -1.48, 3.17),
            ("ITA", 2.89, 0.22, 2.67),
        ],
        decline_rows=[
            ("BNO", -6.29, 8.90, -15.19),
            ("UCO", -6.18, 8.97, -15.15),
            ("USO", -5.48, 7.85, -13.33),
            ("XOP", -6.01, 5.13, -11.14),
            ("XLE", -4.93, 3.48, -8.41),
            ("ICLN", -1.15, 4.80, -5.95),
            ("DBC", -1.82, 3.18, -5.00),
            ("UNG", -0.28, 3.88, -4.16),
            ("XLU", -3.07, 0.80, -3.87),
            ("AMLP", -0.48, 3.31, -3.79),
        ],
        main_candidates=[
            {
                "主標": "AI / Robotics 鍊條深化 — 半導體 reignite 加 人形機器人浮現",
                "核心數據": "SMH +5.94%  ·  IRBO +7.99%  ·  SOXX +5.71%",
                "相關 ETF": "SMH, SOXX, IRBO, THNQ, AIQ, XLK, BOTZ, ROBO, KWEB",
                "driver": "AMD Q1 +38% / TSMC 52w high / Morgan Stanley 5/7 humanoid 報告",
                "對台灣機構意義": "TSMC / HIWIN 在 Humanoid 100 名單、5/22 NVDA 為下一驗證",
            },
            {
                "主標": "油市急殺 — Hormuz 重啟談判加 Iran 提案",
                "核心數據": "USO -5.48% (YTD 跌穿百分百)  ·  BNO -6.29%  ·  XOP -6.01%",
                "相關 ETF": "USO, BNO, UCO, XOP, XLE, AMLP",
                "driver": "Trump 提案 + Iran ceasefire + Brent $115→$103.54 (-5.8%)",
                "對台灣機構意義": "油價回落減輕通膨壓力、Fed 降息預期可能回升",
            },
            {
                "主標": "黃金 V 字反彈 — 油跌帶實質殖利率回落",
                "核心數據": "GDX +5.29% (前週 -7.66%)  ·  GLD +2.01%  ·  SLV +4.85%",
                "相關 ETF": "GDX, GLD, SLV",
                "driver": "油跌 → 通膨預期降 → 實質殖利率回落 → 黃金 carry 優勢",
                "對台灣機構意義": "0504 期黃金崩盤完美反向、商品內部分化主軸不變",
            },
        ],
        sub_candidates=[
            {
                "主標": "新興市場補漲 — EEM/KWEB 美元疲弱受惠",
                "核心數據": "EEM +3.84%  ·  KWEB +2.57%  ·  FXI +1.03%",
                "相關 ETF": "EEM, KWEB, MCHI, FXI, INDA",
                "driver": "美元疲軟 + 油價跌 + MS 報告強調中國 humanoid VC 占 46%",
                "對台灣機構意義": "新興市場 ETF 配置時點可能成熟",
            },
            {
                "主標": "Biotech / 醫療補漲 — XBI/IBB 落後類股回神",
                "核心數據": "XBI +2.48% (前週 -2.15%)  ·  IBB +1.69% (前週 -1.48%)",
                "相關 ETF": "XBI, IBB, XLV",
                "driver": "落後類股輪動、AI / 半導體 spill-over",
                "對台灣機構意義": "防禦型生技配置時機觀察",
            },
            {
                "主標": "Stagflation 訊號 — ISM Prices 70.7 黏住",
                "核心數據": "ISM Services 53.6 / New Orders -7.1 / Prices 70.7",
                "相關 ETF": "TLT, IEF, TIP, XLU, XLP",
                "driver": "Fed 6 月降息機率 5.1% (CME)、stagflation confirmed",
                "對台灣機構意義": "TLT 配置繼續觀望、防禦類股 XLU 反弱訊號",
            },
            {
                "主標": "國防補漲 — ITA + 地緣議題延續",
                "核心數據": "ITA +2.89% (前週 +0.22%)",
                "相關 ETF": "ITA, XAR",
                "driver": "Hormuz 雖緩解但長線地緣議題仍在",
                "對台灣機構意義": "國防 ETF 加碼考量",
            },
            {
                "主標": "公用 / 防禦反弱 — XLU/XLP 與市場分化",
                "核心數據": "XLU -3.07% (前週 +0.80%)  ·  XLP / XLV 弱",
                "相關 ETF": "XLU, XLP, XLV, RSP",
                "driver": "Narrow leadership 集中半導體 + AI、其他類股弱",
                "對台灣機構意義": "等權 RSP 可能跑輸 SPY、配置需檢視 cap-weight bias",
            },
        ],
        theme_tables=[
            ("AI / 半導體 / Robotics", [
                ["SMH", "VanEck 半導體 (含 TSM)", "+5.94%", "+0.67%", "+48.29%", "0511 主議題核心"],
                ["SOXX", "iShares 半導體純美股", "+5.71%", "+0.90%", "+61.36%", "YTD 創新高"],
                ["IRBO", "iShares Robotics & AI", "+7.99%", "+1.19%", "+39.53%", "本週最強"],
                ["THNQ", "ROBO Global AI", "+6.65%", "+1.77%", "+22.26%", "AI 純"],
                ["AIQ", "Global X AI", "+5.58%", "+3.05%", "+16.61%", "AI 純廣"],
                ["XLK", "科技類股", "+4.83%", "+1.03%", "+16.47%", "大型科技"],
                ["BOTZ", "Global X Robotics", "+3.19%", "+3.32%", "+8.85%", "工業機器人"],
                ["ROBO", "ROBO Global Robotics", "+2.28%", "+2.44%", "+20.45%", "全 Robotics"],
                ["KWEB", "中國互聯網", "+2.57%", "-0.17%", "-14.36%", "中國 humanoid 投資鍊"],
            ]),
            ("商品 / 避險 / 大盤", [
                ["USO", "WTI 期貨", "-5.48%", "+7.85%", "+93.89%", "0504 主議題反向"],
                ["BNO", "Brent 期貨", "-6.29%", "+8.90%", "+88.51%", "對沖中東反向"],
                ["XOP", "油氣 E&P", "-6.01%", "+5.13%", "+31.84%", "高 Beta"],
                ["XLE", "綜合能源", "-4.93%", "+3.48%", "+26.21%", "長線部位仍佳"],
                ["GDX", "黃金礦業", "+5.29%", "-7.66%", "+6.84%", "V 字反彈"],
                ["GLD", "黃金實體", "+2.01%", "-2.32%", "+8.30%", "實質殖利率回落"],
                ["SLV", "白銀", "+4.85%", "-0.73%", "+8.47%", "工業屬性 + 避險"],
                ["TLT", "20+ 年公債", "+0.05%", "-0.90%", "-1.30%", "Fed 預期 stable"],
                ["SPY", "標普 500", "+1.52%", "+0.94%", "+6.65%", "5/7 創 YTD 新高"],
            ]),
        ],
        macro_events=[
            "Hormuz 重啟談判 (Trump 提案 + Iran ceasefire 維持) — Brent $115→$103.54 / -5.8%",
            "AMD Q1 2026 營收 $10.3B (+38% YoY) — ignited semiconductor rally",
            "TSMC 5/7 +6% 至 $418 — 52-week high",
            "ISM Services PMI 4 月 53.6 (連 22 月擴張) — Prices 70.7 黏住、stagflation confirmed",
            "Fed 6 月降息機率 5.1% (CME FedWatch 5/7) — 94.9% 預期 hold 3.50-3.75%",
            "5/8 NFP (ET 8:30am) — Monday patch 用",
            "Gold spot 5/7 $4,705 — 自 4/28 $4,620 反彈 +1.8%",
        ],
        research_reports=[
            "Morgan Stanley 5/7 「Humanoids: Humanoid Horizons: Money Meets Machines」 — Tech Diffusion 2026 Key Theme",
            "  · 2026 YTD VC 融資已超 2025 全年 (China 占 46%, April 41 deals)",
            "  · Bezos Project Prometheus $38B / Meta x ARI / Unitree IPO",
            "  · Humanoid 100 Index 自 2025/2/6 +45% (跑贏 SPX/MSCI)",
        ],
        monday_tasks=[
            "議題組合確認：1 主 + 2 補充（從 8 候選圈選）",
            "重抓 yfinance 含 5/8 close — 5 trading days 完整週",
            "NLM saturation Round 1+ 加 Codex Gate 3 平行驗證",
            "寫 reading_map.md + 底稿.md（UWS 12 章節）",
            "build docx + lint_ryan_feedback 12 規則",
            "NLM Studio Path A 全來源 + Path B 只 docx + Codex Path C deterministic",
            "Codex audit Path A/B/C 一致性 + Codex Ryan persona pre-screen",
            "quality_report A++ 簽核 + reviewed.docx canonical",
        ],
        out_path=out_path,
    )
